#!/usr/bin/env python3
"""Generate slides comparing unified analysis + profiling feedback results."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(0, 70, 130)
BLACK = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
GREEN = RGBColor(20, 120, 20)
RED = RGBColor(180, 30, 30)
ORANGE = RGBColor(180, 90, 0)
LIGHT_BLUE = RGBColor(220, 235, 250)
LIGHT_GRAY = RGBColor(242, 242, 247)
LIGHT_GREEN = RGBColor(220, 245, 220)
LIGHT_RED = RGBColor(250, 225, 225)
LIGHT_ORANGE = RGBColor(255, 235, 210)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def title_bar(slide, title, subtitle=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    s.fill.solid(); s.fill.fore_color.rgb = LIGHT_BLUE; s.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.5))
    p = tx.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK_BLUE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.4))
        p2 = tx2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(14); p2.font.color.rgb = GRAY


def add_table(slide, data, left, top, width, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.33 * rows))
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = Inches(w)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK if r > 0 else DARK_BLUE
                if r == 0: p.font.bold = True
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BLUE
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY


def info_box(slide, left, top, w, h, border, bg, title, lines, tsz=16, lsz=13):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.color.rgb = border; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(tsz); p.font.bold = True; p.font.color.rgb = border
    for t in lines:
        p = tf.add_paragraph(); p.text = t; p.font.size = Pt(lsz); p.font.color.rgb = BLACK


def text_block(slide, left, top, w, h, lines, fsz=14):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, (txt, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(fsz); p.font.bold = bold; p.font.color.rgb = color or BLACK


def parse_log(path):
    """Extract per-kernel initial/final speedup and improvements from log."""
    data = {}
    current = None
    with open(path) as f:
        for line in f:
            s = line.strip()
            if s.startswith("Processing:"):
                current = s.split("Processing:")[1].strip()
                data[current] = {"initial": None, "final": None, "improved": False, "ncu": {}}
            elif "Benchmark:" in s and "x speedup" in s and current:
                sp = float(s.split("Benchmark:")[1].split("x")[0].strip())
                if data[current]["initial"] is None:
                    data[current]["initial"] = sp
                    data[current]["final"] = sp
            elif "IMPROVEMENT:" in s and current:
                after = float(s.split("->")[1].strip().rstrip("x"))
                data[current]["final"] = after
                data[current]["improved"] = True
            elif "Bottleneck:" in s and current:
                data[current]["ncu"]["bottleneck"] = s.split("Bottleneck:")[1].strip()
            elif "SM:" in s and "Mem:" in s and current:
                for part in s.split(","):
                    part = part.strip()
                    if part.startswith("SM:"): data[current]["ncu"]["sm"] = part.split(":")[1].strip().rstrip("%")
                    elif part.startswith("Mem:"): data[current]["ncu"]["mem"] = part.split(":")[1].strip().rstrip("%")
                    elif part.startswith("Occ:"): data[current]["ncu"]["occ"] = part.split(":")[1].strip().rstrip("%")
                    elif part.startswith("L1 hit:"): data[current]["ncu"]["l1"] = part.split(":")[1].strip().rstrip("%")
    return data


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), prs.slide_width, Inches(3.5))
s.fill.solid(); s.fill.fore_color.rgb = LIGHT_BLUE; s.line.fill.background()
tx = slide.shapes.add_textbox(Inches(1), Inches(1.9), Inches(11), Inches(1.5))
tf = tx.text_frame
p = tf.paragraphs[0]; p.text = "Unified Analysis + Profiling Feedback"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = DARK_BLUE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "From Pattern-Specific Prompts to General-Purpose Analysis"; p2.font.size = Pt(20); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph(); p3.text = ""; p3.font.size = Pt(10)
p4 = tf.add_paragraph(); p4.text = "PolyBench/C (30 kernels) | 1x Scale"; p4.font.size = Pt(18); p4.font.color.rgb = GRAY; p4.alignment = PP_ALIGN.CENTER
tx2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.8))
p5 = tx2.text_frame.paragraphs[0]; p5.text = "Xiao Qin | University of Leeds | April 2026"; p5.font.size = Pt(18); p5.font.color.rgb = GRAY; p5.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 2: Architecture Before/After
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Architecture: Before vs After")

info_box(slide, 0.4, 1.4, 5.8, 3.0, RED, LIGHT_RED,
    "Before: Pattern-Specific Prompts (870 lines)", [
    "",
    "16 analysis modules, each with custom output format",
    "build_polybench_prompt(): 870 lines of if/elif/else",
    "  - Special case for triangular, stencil, multi-phase...",
    "  - Prescriptive code templates per pattern",
    "  - Adding new pattern = modify prompt builder",
    "",
    "Fragile, hard to maintain, doesn't scale",
])

info_box(slide, 6.8, 1.4, 6.0, 3.0, GREEN, LIGHT_GREEN,
    "After: Unified Analysis + Profiling (250 lines)", [
    "",
    "kernel_analysis.py: single analyze_kernel() function",
    "  -> Runs all analysis passes",
    "  -> Returns structured JSON (fixed schema)",
    "format_analysis_for_prompt(): pattern-agnostic renderer",
    "  -> 80 lines, renders whatever the analysis contains",
    "  -> Adding new pattern = add to JSON schema",
    "",
    "Clean, general, scalable",
])

info_box(slide, 0.4, 4.7, 12.4, 2.3, ORANGE, LIGHT_ORANGE,
    "Key Design Principle", [
    "",
    "Old: Tell the LLM exactly WHAT TO DO for each pattern (prescriptive templates)",
    "New: Tell the LLM WHAT IS TRUE about the kernel (analysis facts) and let it reason",
    "",
    "The compiler analysis guarantees WHAT is safe. The LLM decides HOW.",
    "Profiling feedback then iteratively improves the HOW.",
])

# ============================================================
# SLIDE 3: 4-Way Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "4-Way Comparison: PolyBench/C 1x Scale (30 kernels)")

data = [
    ["Configuration", "Pass Rate", "1st Try", "Median", "Mean", "Max", ">1x", "Lines of Prompt Logic"],
    ["No analysis (baseline)", "28/30 (93%)", "—", "1.06x", "1.52x", "—", "14/28", "0"],
    ["Legacy (pattern-specific)", "30/30 (100%)", "21", "1.77x", "2.37x", "10.01x", "19/30", "870"],
    ["Unified analysis", "30/30 (100%)", "18", "1.40x", "1.90x", "10.17x", "16/30", "80"],
    ["Unified + profiling", "30/30 (100%)", "15", "1.90x", "2.36x", "10.34x", "19/29", "80 + NCU loop"],
]
add_table(slide, data, 0.3, 1.3, 12.7, col_widths=[2.8, 1.2, 0.8, 1.0, 1.0, 1.0, 1.0, 2.0])

info_box(slide, 0.4, 4.2, 12.4, 2.8, DARK_BLUE, LIGHT_BLUE,
    "Interpretation", [
    "",
    "1. Unified analysis alone: slightly lower performance (1.40x vs 1.77x median)",
    "   — expected: general guidance vs prescriptive templates",
    "",
    "2. Unified + profiling feedback: MATCHES legacy (1.90x vs 1.77x median, 2.36x vs 2.37x mean)",
    "   — profiling feedback closes the gap while maintaining generality",
    "",
    "3. Same 100% correctness across all analysis-guided configurations",
    "   — compiler analysis is the correctness foundation regardless of prompt style",
])

# ============================================================
# SLIDE 4: Profiling Improvements Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Profiling Feedback: 15 Kernels Improved")

# Parse the unified+profiling log
prof_data = parse_log("/tmp/unified_profiling_1x.log")
improved = [(k, d) for k, d in sorted(prof_data.items()) if d["improved"]]

imp_table = [["Kernel", "Before", "After", "Improvement", "Bottleneck"]]
for k, d in sorted(improved, key=lambda x: -(x[1]["final"] / x[1]["initial"] if x[1]["initial"] else 0)):
    init = d["initial"]
    final = d["final"]
    if init and init > 0:
        pct = f"+{((final/init)-1)*100:.0f}%"
    else:
        pct = "—"
    bn = d["ncu"].get("bottleneck", "—")
    if "memory" in bn: bn = "memory-bound"
    elif "compute" in bn: bn = "compute-bound"
    elif "latency" in bn: bn = "latency-bound"
    else: bn = bn[:15]
    imp_table.append([k, f"{init:.2f}x", f"{final:.2f}x", pct, bn])

add_table(slide, imp_table, 0.3, 1.3, 12.7, col_widths=[2.0, 1.5, 1.5, 1.5, 4.0])

# ============================================================
# SLIDE 5: NCU Profiling Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "NCU Profiling Metrics for All Kernels")

ncu_table = [["Kernel", "Speedup", "SM%", "Mem%", "Occ%", "L1%", "Bottleneck", "Prof. Improved?"]]
for k in sorted(prof_data.keys()):
    d = prof_data[k]
    sp = f"{d['final']:.2f}x" if d["final"] else "—"
    ncu = d["ncu"]
    bn = ncu.get("bottleneck", "—")
    if "memory" in bn: bn = "mem"
    elif "compute" in bn: bn = "compute"
    elif "latency" in bn: bn = "latency"
    else: bn = bn[:8]
    imp = "Y" if d["improved"] else ""
    ncu_table.append([k, sp, ncu.get("sm", "—"), ncu.get("mem", "—"),
                      ncu.get("occ", "—"), ncu.get("l1", "—"), bn, imp])

# Split into two slides if needed
if len(ncu_table) <= 17:
    add_table(slide, ncu_table, 0.3, 1.3, 12.7, col_widths=[1.8, 1.0, 0.8, 0.8, 0.8, 0.8, 1.2, 1.2])
else:
    add_table(slide, ncu_table[:17], 0.3, 1.3, 12.7, col_widths=[1.8, 1.0, 0.8, 0.8, 0.8, 0.8, 1.2, 1.2])
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide2, "NCU Profiling Metrics (continued)")
    add_table(slide2, [ncu_table[0]] + ncu_table[17:], 0.3, 1.3, 12.7,
              col_widths=[1.8, 1.0, 0.8, 0.8, 0.8, 0.8, 1.2, 1.2])

# ============================================================
# SLIDE 6: Key Findings & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Key Findings & Next Steps")

text_block(slide, 0.5, 1.4, 5.8, 5.5, [
    ("Findings", True, DARK_BLUE),
    ("", False, None),
    ("1. General analysis + profiling = pattern-specific hand-tuning", True, GREEN),
    ("   Unified (1.90x) matches legacy (1.77x) median", False, None),
    ("   With 10x less prompt logic (80 vs 870 lines)", False, GRAY),
    ("", False, None),
    ("2. Profiling feedback is most effective on latency-bound kernels", True, DARK_BLUE),
    ("   deriche: 0.12x -> 5.69x, gemver: 0.15x -> 1.98x", False, None),
    ("   LLM identifies insufficient parallelism from NCU metrics", False, GRAY),
    ("", False, None),
    ("3. Correctness remains the bottleneck for optimization", True, RED),
    ("   LLM frequently breaks correctness when optimizing", False, None),
    ("   Our test-based check catches this, but can't fix it", False, GRAY),
    ("", False, None),
    ("4. 100% correctness maintained in ALL analysis-guided configs", True, GREEN),
    ("   Compiler analysis is the correctness foundation", False, None),
], fsz=14)

info_box(slide, 6.8, 1.4, 6.0, 5.5, DARK_BLUE, LIGHT_BLUE,
    "Conference Paper: Combining Both Signals", [
    "",
    "Our system now provides two complementary signals:",
    "",
    "  Compiler analysis -> WHAT is safe (correctness)",
    "  NCU profiling     -> HOW to optimize (performance)",
    "",
    "These are exactly the signals that RL/search need:",
    "",
    "  1. Analysis as correctness constraint/reward",
    "     - Reject candidates that violate dependences",
    "",
    "  2. NCU metrics as performance reward",
    "     - Guide search toward better hardware utilization",
    "",
    "  3. Unified analysis as structured world model",
    "     - K-Search style: decouple strategy from code",
    "",
    "Next: integrate into a search/RL framework",
    "that explores many candidates using both signals",
])

# Save
output = Path("/home/qinxiao/workspace/compiler-guided-triton-gen/presentation/comparison_slides.pptx")
prs.save(str(output))
print(f"Saved to {output}")
