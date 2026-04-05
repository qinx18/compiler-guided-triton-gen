#!/usr/bin/env python3
"""Generate literature review slides for supervision meeting.
All text is dark on light backgrounds — no white text anywhere."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors — all text will be dark; backgrounds are light tints
DARK_BLUE = RGBColor(0, 70, 130)
BLACK = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
LIGHT_BLUE_BG = RGBColor(220, 235, 250)
LIGHT_GRAY_BG = RGBColor(242, 242, 247)

RL_COLOR = RGBColor(20, 80, 140)
RL_BG = RGBColor(215, 230, 250)
SEARCH_COLOR = RGBColor(180, 90, 0)
SEARCH_BG = RGBColor(255, 235, 210)
MEM_COLOR = RGBColor(20, 110, 20)
MEM_BG = RGBColor(215, 240, 215)
NPU_COLOR = RGBColor(180, 30, 30)
NPU_BG = RGBColor(250, 220, 220)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_bar(slide, title_text, subtitle_text=None):
    """Title bar: light blue background, dark text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_body_text(slide, left, top, width, height, bullets, font_size=16, spacing=Pt(6)):
    """Add bulleted text — all dark colors."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (level, text, bold, color) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(font_size - level * 2)
        p.font.bold = bold
        p.font.color.rgb = color if color else BLACK
        p.space_after = spacing
        if level == 0:
            p.space_before = Pt(4)
        if level > 0:
            p.space_after = Pt(2)
    return txBox


def add_tag_box(slide, left, top, width, height, bg_color, text_color, text, font_size=13):
    """Small colored tag/label box — light bg, dark text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_info_box(slide, left, top, width, height, border_color, bg_color, title, lines, title_size=16, line_size=13):
    """Info box with border, light bg, dark text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = border_color
    for text in lines:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(line_size)
        p.font.color.rgb = BLACK
        if any(marker in text for marker in ["→", "×", "31.6%", "47.8%", "97.4%", "2.12", "14.3", "1030", "98.1%", "90.4%", "95.5%"]):
            p.font.bold = True
    return shape


def add_takeaway_box(slide, left, top, width, height, border_color, bg_color, text):
    """Takeaway box — tinted bg with dark text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = border_color
    p.font.bold = True
    return shape


def add_table_slide(slide, data, left, top, width, col_widths=None):
    """Table — header row is light blue bg with dark text."""
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.4 * rows))
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = BLACK if r > 0 else DARK_BLUE
                if r == 0:
                    paragraph.font.bold = True
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY_BG
    return ts


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), prs.slide_width, Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE_BG
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Literature Review: LLM/Agent-Based Kernel Generation"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11), Inches(0.8))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "9 Recent Works (Jan -- Mar 2026)"
p2.font.size = Pt(22)
p2.font.color.rgb = GRAY
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.8))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Xiao Qin  |  Supervision Meeting  |  April 2026"
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 2: Overview Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Overview: 9 Recent Works")

data = [
    ["Paper", "Target", "Base Method", "Key Innovation", "Affiliation", "Date"],
    ["Dr. Kernel", "Triton", "RL fine-tuning", "Multi-turn RL + profiling rewards", "HKUST", "Feb 2026"],
    ["Fine-Tuning GPT-5", "Triton", "RL fine-tuning", "RL post-training of frontier model", "Tehrani et al.", "Feb 2026"],
    ["DICE", "CUDA", "RL fine-tuning", "Diffusion (non-AR) architecture", "Bai et al.", "Feb 2026"],
    ["AVO", "CUDA/PTX", "Evolutionary search", "LLM agents as variation operators", "NVIDIA", "Mar 2026"],
    ["K-Search", "Triton", "Evolutionary search", "Co-evolving world model for planning", "UC Berkeley", "Feb 2026"],
    ["OptiML", "CUDA", "MCTS search", "Nsight Compute profiler in the loop", "Bhattacharjee et al.", "Feb 2026"],
    ["KernelBlaster", "CUDA", "Iterative search", "Persistent cross-task memory", "NVIDIA / Berkeley", "Feb 2026"],
    ["AscendKernelGen", "Ascend NPU", "SFT + RL", "CoT dataset for NPU reasoning", "Cao et al.", "Jan 2026"],
    ["AscendCraft", "Ascend NPU", "Prompting", "Intermediate DSL + LLM lowering", "Wen et al.", "Jan 2026"],
]
add_table_slide(slide, data, 0.3, 1.4, 12.7, col_widths=[1.9, 1.2, 1.7, 3.0, 2.2, 1.0])

for i, (label, bg, tc) in enumerate([
    ("RL-Based", RL_BG, RL_COLOR), ("Search/Evolution", SEARCH_BG, SEARCH_COLOR),
    ("Memory-Augmented", MEM_BG, MEM_COLOR), ("NPU-Targeted", NPU_BG, NPU_COLOR)
]):
    add_tag_box(slide, 1.0 + i * 3.0, 6.6, 2.2, 0.4, bg, tc, label, font_size=12)

# ============================================================
# SLIDE 3: What Do These Categories Mean?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "What Do These Categories Mean?")

# SFT vs RL explanation
add_info_box(slide, 0.4, 1.4, 12.4, 1.6, DARK_BLUE, LIGHT_BLUE_BG,
    "Background: Two ways to fine-tune an LLM", [
    "",
    "SFT (Supervised Fine-Tuning): Give the model (input, correct output) pairs, train it to mimic. Simple but needs high-quality data.",
    "RL Fine-Tuning: Let the model generate code, compile & run it, give a reward score. Model learns from trial-and-error. Harder to train but doesn't need 'correct' answers.",
    "Typical pipeline: SFT first (warm-up) -> then RL on top (push beyond what examples can teach).",
])

# RL Fine-Tuning explanation
add_info_box(slide, 0.4, 3.3, 6.0, 1.8, RL_COLOR, RL_BG,
    "RL Fine-Tuning (modify the LLM's weights)", [
    "",
    "Further train the LLM on kernel tasks using rewards",
    "from compiling & running the generated code.",
    "Like grading homework repeatedly until the student improves.",
])

# Search/Evolution explanation
add_info_box(slide, 6.8, 3.3, 6.0, 1.8, SEARCH_COLOR, SEARCH_BG,
    "Search / Evolution (keep LLM frozen, search at inference)", [
    "",
    "Don't change the LLM. Use it to generate many candidates,",
    "test them, and iteratively improve via search (evolutionary",
    "algorithms, MCTS tree search). Like brainstorm + trial-and-error.",
])

# Memory-Augmented explanation
add_info_box(slide, 0.4, 5.4, 6.0, 1.8, MEM_COLOR, MEM_BG,
    "Memory-Augmented (accumulate experience, no weight updates)", [
    "",
    "LLM is frozen, but the system stores what optimizations",
    "worked before. Retrieves relevant past experience as prompt",
    "context. Like an engineer with lab notebooks.",
])

# NPU-Targeted explanation
add_info_box(slide, 6.8, 5.4, 6.0, 1.8, NPU_COLOR, NPU_BG,
    "NPU-Targeted (non-GPU accelerators)", [
    "",
    "LLMs know nearly nothing about NPU programming (Huawei",
    "Ascend) -- need domain-specific datasets, DSLs, or CoT.",
    "'DSL transcompilation' = LLM writes simpler DSL, then lower to HW.",
])

# ============================================================
# SLIDE 4: Taxonomy Diagram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Taxonomy of Approaches")

add_tag_box(slide, 4.5, 1.6, 4.3, 0.6, LIGHT_BLUE_BG, DARK_BLUE, "LLM-Based Kernel Generation", 16)

cats = [
    (0.8, "RL Fine-Tuning", RL_BG, RL_COLOR, ["Dr. Kernel (14B, TRLOO)", "FT GPT-5 (RL post-train)", "DICE (Diffusion + BiC-RL)"]),
    (3.8, "Search / Evolution", SEARCH_BG, SEARCH_COLOR, ["AVO (agent evolutionary ops)", "K-Search (world model)", "OptiML (MCTS + Nsight)"]),
    (7.0, "Memory-Augmented", MEM_BG, MEM_COLOR, ["KernelBlaster (MAIC-RL)"]),
    (9.8, "NPU-Targeted", NPU_BG, NPU_COLOR, ["AscendKernelGen (SFT+RL)", "AscendCraft (DSL transc.)"]),
]

for x, label, bg, tc, papers in cats:
    add_tag_box(slide, x, 2.8, 2.8, 0.5, bg, tc, label, 14)
    for j, paper in enumerate(papers):
        txBox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.6 + j * 0.45), Inches(2.6), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "- " + paper
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK

# Insight box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY_BG
shape.line.color.rgb = DARK_BLUE
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(12)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "Two dominant paradigms:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
for t in [
    "1) Train the model (RL fine-tuning) -- modify weights to internalize kernel optimization knowledge",
    "2) Search at inference time (evolution, MCTS, agent loops) -- keep model frozen, explore solution space with feedback",
    "3) NPU works address the cold-start problem: LLMs have near-zero NPU knowledge, requiring domain-specific datasets/DSLs",
]:
    p = tf.add_paragraph()
    p.text = t
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK

# ============================================================
# SLIDE 4: Dr. Kernel
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Dr. Kernel: RL Done Right for Triton Kernels", "Liu et al. (HKUST) | arXiv 2602.05885 | Feb 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.5, [
    (0, "Problem", True, DARK_BLUE),
    (1, "Multi-turn RL for kernel generation suffers from reward hacking", False, None),
    (1, "Standard REINFORCE has biased gradients in multi-turn settings", False, None),
    (0, "Three Key Contributions", True, DARK_BLUE),
    (1, "KernelGYM: distributed GPU env with reward-hacking detection, multi-turn data collection, stable long-term RL training", False, None),
    (1, "TRLOO (Turn-level REINFORCE-Leave-One-Out): unbiased advantage estimation for multi-turn RL", False, None),
    (1, "Profiling-based rewards (PR): uses GPU profiling metrics, not wallclock time -- prevents 'lazy optimization' shortcuts", False, None),
    (1, "Profiling-based rejection sampling (PRS): filters training data by profiling quality", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 3.0, RL_COLOR, RL_BG,
    "Results (KernelBench Level-2)", [
    "",
    "Fraction of kernels >=1.2x speedup over PyTorch:",
    "  Dr. Kernel-14B:  31.6%",
    "  Claude Sonnet 4.5:  26.7%",
    "  GPT-5:  28.6%",
    "",
    "With test-time scaling (best-of-N): 47.8%",
])

add_takeaway_box(slide, 7.0, 4.8, 5.8, 1.0, RL_COLOR, RL_BG,
    "Takeaway: Naive RL rewards (pass/fail, wallclock) are easily hacked. GPU profiling metrics are essential for honest optimization.")

# ============================================================
# SLIDE 5: Fine-Tuning GPT-5
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Fine-Tuning GPT-5 for GPU Kernel Generation", "Tehrani et al. | arXiv 2602.11000 | Feb 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Problem", True, DARK_BLUE),
    (1, "Even frontier LLMs struggle with GPU kernels due to:", False, None),
    (1, "  - Scarcity of high-quality labeled training data", False, GRAY),
    (1, "  - Compiler biases when generating synthetic solutions", False, GRAY),
    (1, "  - Limited generalization across hardware generations", False, GRAY),
    (0, "Method", True, DARK_BLUE),
    (1, "RL fine-tuning of GPT-5 using Makora's environment", False, None),
    (1, "RL as a data-efficient alternative to SFT", False, None),
    (1, "Integrated into a full coding agent for multi-turn refinement", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 3.5, RL_COLOR, RL_BG,
    "Results", [
    "",
    "Single-attempt correctness:",
    "  43.7% → 77.0% (+33.3pp)",
    "",
    "Kernels outperforming TorchInductor:",
    "  14.8% → 21.8% (+7pp)",
    "",
    "With full coding agent:",
    "  97.4% solve rate",
    "  2.12× geomean speedup over TorchInductor",
])

add_takeaway_box(slide, 7.0, 5.3, 5.8, 0.8, RL_COLOR, RL_BG,
    "Takeaway: Even frontier models have huge headroom for domain-specific RL post-training. Agent wrapper pushes correctness to near-100%.")

# ============================================================
# SLIDE 6: DICE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "DICE: Diffusion LLMs for CUDA Kernels", "Bai et al. | arXiv 2602.11715 | Feb 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Key Idea: Use diffusion LLMs (not autoregressive)", True, DARK_BLUE),
    (0, "Why Diffusion?", True, DARK_BLUE),
    (1, "Parallel token generation -> holistic structural planning", False, None),
    (1, "Non-sequential refinement suits code where global structure matters", False, None),
    (1, "Avoids left-to-right commitment of autoregressive models", False, None),
    (0, "Method", True, DARK_BLUE),
    (1, "CuKe dataset: curated SFT data, 2.0x speedup threshold (quality > quantity)", False, None),
    (1, "BiC-RL (Bi-phase Curated RL):", False, None),
    (1, "  Stage 1: Kernel infilling (complete partial kernels -- easier task)", False, GRAY),
    (1, "  Stage 2: End-to-end generation (full kernels from scratch)", False, GRAY),
    (1, "Models at 1.7B, 4B, 8B parameter scales", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 2.5, RL_COLOR, RL_BG,
    "Results (KernelBench)", [
    "",
    "Outperforms both AR and diffusion LLMs",
    "of comparable scale on KernelBench",
    "",
    "First diffusion LLM family specifically",
    "designed for kernel generation",
])

add_takeaway_box(slide, 7.0, 4.3, 5.8, 0.8, RL_COLOR, RL_BG,
    "Takeaway: Diffusion architectures are a viable alternative to AR for structured code generation -- holistic planning matters.")

# ============================================================
# SLIDE 7: AVO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "AVO: Agentic Variation Operators", "Chen et al. (NVIDIA) | arXiv 2603.24517 | Mar 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Key Idea: Replace evolutionary operators with LLM agents", True, DARK_BLUE),
    (0, "How It Works", True, DARK_BLUE),
    (1, "Classical: fixed mutation + crossover heuristics", False, GRAY),
    (1, "AVO: self-directed agent loops as variation operators", False, None),
    (1, "Each agent can:", False, None),
    (1, "  - Consult evolutionary lineage (what worked/failed)", False, GRAY),
    (1, "  - Query a domain-specific knowledge base", False, GRAY),
    (1, "  - Use execution feedback to propose, repair, critique, verify", False, GRAY),
    (1, "7 days of continuous autonomous evolution on Blackwell B200", False, None),
    (0, "Target: Multi-head Attention", True, DARK_BLUE),
    (1, "The most aggressively hand-optimized kernel in AI", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 2.5, SEARCH_COLOR, SEARCH_BG,
    "Results (Attention on B200)", [
    "",
    "Outperforms cuDNN by up to 3.5%",
    "Outperforms FlashAttention-4 by up to 10.5%",
    "",
    "Optimizations transfer across",
    "attention variants",
])

add_takeaway_box(slide, 7.0, 4.3, 5.8, 1.0, SEARCH_COLOR, SEARCH_BG,
    "Takeaway: LLM agents can beat hand-tuned vendor libraries on the most optimized kernels -- but requires massive compute (7 days on B200s).")

# ============================================================
# SLIDE 8: K-Search
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "K-Search: Co-Evolving World Model", "Cao et al. (UC Berkeley) | arXiv 2602.19128 | Feb 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.5, [
    (0, "Problem", True, DARK_BLUE),
    (1, "Existing evolutionary approaches treat LLMs as stochastic code generators", False, None),
    (1, "Lack planning; discard promising strategies due to buggy intermediate code", False, None),
    (0, "Key Innovation: Decoupled Planning & Instantiation", True, DARK_BLUE),
    (1, "High-level: what optimization strategy to pursue", False, None),
    (1, "  (e.g., 'use shared memory tiling with 2D thread blocks')", False, GRAY),
    (1, "Low-level: actual code implementation", False, None),
    (1, "Co-evolving world model guides search using LLM domain knowledge", False, None),
    (1, "Enables non-monotonic optimization paths:", False, None),
    (1, "  temporary regressions -> large gains (won't abandon buggy but promising strategies)", False, GRAY),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 3.0, SEARCH_COLOR, SEARCH_BG,
    "Results", [
    "",
    "Average 2.10× over SOTA evolutionary search",
    "Up to 14.3× on complex MoE kernels",
    "",
    "SOTA on GPUMode TriMul benchmark:",
    "  1030us on H100",
    "  Beats both automated & human solutions",
])

add_takeaway_box(slide, 7.0, 4.8, 5.8, 0.9, SEARCH_COLOR, SEARCH_BG,
    "Takeaway: Separating 'what to try' from 'how to implement it' is crucial -- don't abandon strategies just because intermediate code is buggy.")

# ============================================================
# SLIDE 9: OptiML
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "OptiML: MCTS + Hardware Profiler Feedback", "Bhattacharjee et al. | arXiv 2602.12305 | Feb 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Key Idea: Kernel optimization as search under verification", True, DARK_BLUE),
    (0, "Two-Stage Architecture", True, DARK_BLUE),
    (1, "OptiML-G (Generator): Mixture-of-Thoughts proposal policy", False, None),
    (1, "  - Synthesizes initial programs from NL or existing code", False, GRAY),
    (1, "OptiML-X (Optimizer): Monte Carlo Tree Search", False, None),
    (1, "  - LLM-driven code edits guided by Nsight Compute feedback", False, GRAY),
    (1, "  - Each candidate: compiled -> verified -> profiled", False, GRAY),
    (0, "Reward Design", True, DARK_BLUE),
    (1, "Composite: runtime + HW bottleneck (mem throughput, occupancy) + regression guards", False, None),
    (0, "Key Advantage", True, DARK_BLUE),
    (1, "Interpretable: traceable sequence of edits with profiler evidence", False, None),
    (1, "Stages decoupled: can use OptiML-X on human-written kernels", False, None),
], font_size=15)

add_takeaway_box(slide, 7.0, 5.0, 5.8, 0.8, SEARCH_COLOR, SEARCH_BG,
    "Takeaway: MCTS + Nsight Compute feedback = better performance + explainability. Two decoupled stages are independently useful.")

# ============================================================
# SLIDE 10: KernelBlaster
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "KernelBlaster: Cross-Task Memory-Augmented RL", "Dong et al. (NVIDIA / Berkeley) | arXiv 2602.14293 | Feb 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Problem", True, DARK_BLUE),
    (1, "Existing agentic workflows cannot aggregate knowledge from prior optimization", False, None),
    (1, "Re-explore the same space for each new kernel", False, None),
    (0, "MAIC-RL (Memory-Augmented In-Context RL)", True, DARK_BLUE),
    (1, "Persistent, searchable CUDA knowledge base across tasks", False, None),
    (1, "Insights from kernel A inform optimization of kernel B", False, None),
    (1, "No LLM fine-tuning -- learns through accumulated in-context experience", False, None),
    (0, "Profile-Guided Textual Gradients", True, DARK_BLUE),
    (1, "Profiling feedback encoded as 'gradients' in text form", False, None),
    (1, "Steers exploration toward promising optimization regions", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 2.0, MEM_COLOR, MEM_BG,
    "Results (vs PyTorch baselines)", [
    "",
    "Level 1: 1.43× speedup",
    "Level 2: 2.50× speedup",
    "Level 3: 1.50× speedup",
])

add_takeaway_box(slide, 7.0, 4.0, 5.8, 0.8, MEM_COLOR, MEM_BG,
    "Takeaway: Cross-task knowledge transfer avoids redundant exploration. Fine-tuning not always needed if you can accumulate experience.")

# ============================================================
# SLIDE 11: AscendKernelGen
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "AscendKernelGen: LLMs for NPU Kernels", "Cao et al. (20 authors) | arXiv 2601.07160 | Jan 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Problem: LLMs completely fail on NPU kernels", True, DARK_BLUE),
    (1, "General-purpose LLMs: near-zero success on Ascend NPU", False, None),
    (1, "NPU DSLs are too specialized, underrepresented in training data", False, None),
    (0, "Key Contributions", True, DARK_BLUE),
    (1, "Ascend-CoT dataset: chain-of-thought reasoning for NPU kernel programming", False, None),
    (1, "KernelGen-LM: domain-specific LM with SFT + RL + execution feedback", False, None),
    (1, "NPUKernelBench: evaluation across compilation, correctness, performance", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 2.5, NPU_COLOR, NPU_BG,
    "Results (Complex Level-2 Kernels)", [
    "",
    "Compilation success:",
    "  0% (baseline) -> 95.5% (Pass@10)",
    "",
    "Functional correctness: 64.3%",
])

add_takeaway_box(slide, 7.0, 4.3, 5.8, 0.8, NPU_COLOR, NPU_BG,
    "Takeaway: GPU-centric LLM training data does NOT transfer to NPUs. Domain-specific datasets with structured reasoning are essential.")

# ============================================================
# SLIDE 12: AscendCraft
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "AscendCraft: DSL-Guided Transcompilation for NPU", "Wen et al. | arXiv 2601.22760 | Jan 2026")

add_body_text(slide, 0.6, 1.3, 6.0, 5.0, [
    (0, "Key Idea: Decompose hard problem via intermediate DSL", True, DARK_BLUE),
    (0, "Two-Step Pipeline", True, DARK_BLUE),
    (1, "Step 1 -- DSL Generation:", False, None),
    (1, "  LLM generates lightweight DSL abstracting AscendC complexity", False, GRAY),
    (1, "  Preserves Ascend-specific semantics (tiling, data movement, pipeline)", False, GRAY),
    (1, "  Uses category-specific expert examples (7 operator categories)", False, GRAY),
    (1, "Step 2 -- Transcompilation:", False, None),
    (1, "  Constraint-driven LLM lowering passes: DSL -> full AscendC", False, GRAY),
    (0, "Generalizability", True, DARK_BLUE),
    (1, "Successfully generates correct kernels for mHC architecture (different NPU)", False, None),
], font_size=15)

add_info_box(slide, 7.0, 1.5, 5.8, 2.5, NPU_COLOR, NPU_BG,
    "Results (MultiKernelBench)", [
    "",
    "Compilation: 98.1%",
    "Functional correctness: 90.4%",
    "46.2% match/exceed PyTorch eager perf",
    "",
    "Generalizes to mHC architecture",
])

add_takeaway_box(slide, 7.0, 4.3, 5.8, 0.8, NPU_COLOR, NPU_BG,
    "Takeaway: An intermediate DSL decomposes the hard problem into two tractable steps. The approach is portable across NPU targets.")

# ============================================================
# SLIDE 13: Cross-Paper Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Cross-Paper Comparison")

data = [
    ["Paper", "Fine-tune?", "Multi-turn?", "Profiler?", "Search?", "Open-source?"],
    ["Dr. Kernel", "Yes (14B)", "Yes", "Yes (profiling rewards)", "No", "Yes (KernelGYM)"],
    ["FT GPT-5", "Yes (GPT-5)", "Yes (agent)", "No", "No", "No"],
    ["DICE", "Yes (1.7-8B)", "No", "No", "No", "Partial"],
    ["AVO", "No", "Yes (agent)", "Yes", "Evolutionary", "No"],
    ["K-Search", "No", "Yes", "Yes", "Evol. + world model", "Yes"],
    ["OptiML", "No", "Yes", "Yes (Nsight)", "MCTS", "Unknown"],
    ["KernelBlaster", "No", "Yes (agent)", "Yes", "In-context RL", "Yes"],
    ["AscendKernelGen", "Yes", "No", "No", "No", "Partial"],
    ["AscendCraft", "No", "No", "No", "No", "Unknown"],
]
add_table_slide(slide, data, 0.3, 1.4, 12.7, col_widths=[2.0, 1.5, 1.5, 2.5, 2.5, 1.5])

# ============================================================
# SLIDE 14: Key Trends
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Key Trends Across All 9 Works")

add_body_text(slide, 0.6, 1.4, 5.8, 5.5, [
    (0, "1. RL fine-tuning is the dominant training paradigm", True, DARK_BLUE),
    (1, "4/9 actually update model weights via RL (Dr.Kernel, FT GPT-5, DICE, AscendKernelGen)", False, None),
    (1, "Note: KernelBlaster's 'in-context RL' is NOT weight-updating RL --", False, None),
    (1, "  it keeps the LLM frozen and 'learns' by accumulating experience in prompts", False, GRAY),
    (1, "Profiling-based rewards >> pass/fail or wallclock", False, None),
    (0, "2. Search as alternative to fine-tuning", True, DARK_BLUE),
    (1, "AVO, K-Search, OptiML, KernelBlaster avoid fine-tuning", False, None),
    (1, "Trade training cost for inference-time compute", False, None),
    (0, "3. Decoupling planning from implementation", True, DARK_BLUE),
    (1, "Idea: decide WHAT optimization to try (e.g. 'use shared", False, None),
    (1, "  memory tiling') separately from HOW to write the code", False, GRAY),
    (1, "Why it matters: a good strategy may have a buggy first", False, None),
    (1, "  implementation. Without decoupling, the system sees bad", False, GRAY),
    (1, "  perf and abandons the strategy. With decoupling, it can", False, GRAY),
    (1, "  retry the implementation while keeping the strategy.", False, GRAY),
    (1, "K-Search, OptiML, AscendCraft all use this pattern", False, None),
    (0, "4. Hardware feedback loops are essential", True, DARK_BLUE),
    (1, "5/9 works use profiler feedback (Nsight, GPU profiling)", False, None),
    (0, "5. Non-GPU targets emerging", True, DARK_BLUE),
    (1, "2 works target Ascend NPUs; GPU knowledge doesn't transfer", False, None),
], font_size=14)

add_info_box(slide, 7.0, 1.5, 5.8, 4.0, DARK_BLUE, LIGHT_BLUE_BG,
    "What's notably absent", [
    "",
    "- No work uses compiler analysis to guide",
    "  generation (our unique angle)",
    "",
    "- Few address correctness verification",
    "  beyond test-case matching",
    "",
    "- No work targets C-to-Triton translation",
    "  (most are PyTorch-to-Triton/CUDA)",
    "",
    "- Limited focus on dependence analysis",
    "  for race condition prevention",
])

# ============================================================
# SLIDE 15: Positioning & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Positioning Our Work & Next Steps")

# Left: current work
add_info_box(slide, 0.5, 1.5, 5.8, 2.5, DARK_BLUE, LIGHT_BLUE_BG,
    "Current Work (Workshop Paper)", [
    "",
    "- Compiler analysis -> structured prompts -> LLM -> Triton",
    "- Focus: what parallelism is safe (correctness-first)",
    "- 30/30 PolyBench/C, 150/151 TSVC",
    "",
    "Gap: single-shot prompting only",
    "No RL, no search, no iterative refinement",
])

# Right: extensions
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE_BG
shape.line.color.rgb = DARK_BLUE
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(12)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "Extension Directions (Conference Paper)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
for text in [
    "",
    "1. Analysis-guided search: compiler analysis as",
    "   'world model' in K-Search / MCTS framework",
    "2. Analysis as RL reward: dependence analysis",
    "   for correctness verification in RL loop",
    "3. Multi-turn refinement: analysis feedback to",
    "   iteratively fix race conditions",
    "4. Harder benchmarks: KernelBench-style tasks",
]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    if text.startswith("1.") or text.startswith("2.") or text.startswith("3.") or text.startswith("4."):
        p.font.bold = True

# Discussion questions at bottom
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SEARCH_BG
shape2.line.color.rgb = SEARCH_COLOR
shape2.line.width = Pt(2)
tf2 = shape2.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(12)
tf2.margin_top = Pt(8)
p = tf2.paragraphs[0]
p.text = "Discussion Points for Supervisors"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = SEARCH_COLOR
for text in [
    "",
    "1. Which extension direction best leverages our unique compiler analysis angle?",
    "2. Search-based (K-Search style) vs. RL-based (Dr. Kernel style) -- which is more feasible for us?",
    "3. Benchmark choice: stay with PolyBench/TSVC or move to KernelBench (de facto standard)?",
    "4. Conference target and timeline: what is the minimum viable extension to differentiate?",
]:
    p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    if text.startswith("1.") or text.startswith("2.") or text.startswith("3.") or text.startswith("4."):
        p.font.bold = True

# ============================================================
# Save
# ============================================================
output_path = "/home/qinxiao/workspace/compiler-guided-triton-gen/presentation/lit_review_slides.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
