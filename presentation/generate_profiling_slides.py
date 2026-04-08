#!/usr/bin/env python3
"""Generate slides presenting profiling-feedback optimization results.
Run after both 1x and 8x experiments complete."""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0, 70, 130)
BLACK = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
GREEN = RGBColor(20, 120, 20)
RED = RGBColor(180, 30, 30)
LIGHT_BLUE_BG = RGBColor(220, 235, 250)
LIGHT_GRAY_BG = RGBColor(242, 242, 247)
LIGHT_GREEN_BG = RGBColor(220, 245, 220)
LIGHT_RED_BG = RGBColor(250, 225, 225)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_table(slide, data, left, top, width, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.35 * rows))
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = BLACK if r > 0 else DARK_BLUE
                if r == 0:
                    paragraph.font.bold = True
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY_BG
    return ts


def add_text_box(slide, left, top, width, height, lines, font_size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color if color else BLACK
    return txBox


def add_info_box(slide, left, top, width, height, border_color, bg_color, title, lines, title_size=16, line_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = border_color
    for text in lines:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(line_size)
        p.font.color.rgb = BLACK
    return shape


def load_results(results_path):
    """Load results.json and extract per-kernel data."""
    if not results_path.exists():
        return {}
    with open(results_path) as f:
        return json.load(f)


def extract_profiling_data(log_path):
    """Parse the log file to extract profiling iteration data."""
    if not log_path.exists():
        return {}

    data = {}
    current_kernel = None
    current_initial_speedup = None

    with open(log_path) as f:
        for line in f:
            line = line.strip()
            if line.startswith("Processing:"):
                current_kernel = line.split("Processing:")[1].strip()
                data[current_kernel] = {
                    "initial_speedup": None,
                    "final_speedup": None,
                    "iterations": [],
                    "ncu": {},
                    "improved": False,
                }
            elif "Benchmark:" in line and "x speedup" in line and current_kernel:
                sp = float(line.split("Benchmark:")[1].split("x")[0].strip())
                if data[current_kernel]["initial_speedup"] is None:
                    data[current_kernel]["initial_speedup"] = sp
                    data[current_kernel]["final_speedup"] = sp
            elif "Bottleneck:" in line and current_kernel:
                data[current_kernel]["ncu"]["bottleneck"] = line.split("Bottleneck:")[1].strip()
            elif "SM:" in line and "Mem:" in line and current_kernel:
                parts = line.split(",")
                for part in parts:
                    part = part.strip()
                    if part.startswith("SM:"):
                        data[current_kernel]["ncu"]["sm"] = part.split(":")[1].strip().rstrip("%")
                    elif part.startswith("Mem:"):
                        data[current_kernel]["ncu"]["mem"] = part.split(":")[1].strip().rstrip("%")
                    elif part.startswith("Occ:"):
                        data[current_kernel]["ncu"]["occ"] = part.split(":")[1].strip().rstrip("%")
                    elif part.startswith("L1 hit:"):
                        data[current_kernel]["ncu"]["l1"] = part.split(":")[1].strip().rstrip("%")
            elif "IMPROVEMENT:" in line and current_kernel:
                parts = line.split("IMPROVEMENT:")[1].strip()
                before, after = parts.split("->")
                new_sp = float(after.strip().rstrip("x"))
                data[current_kernel]["final_speedup"] = new_sp
                data[current_kernel]["improved"] = True
            elif "Optimized speedup:" in line and current_kernel:
                sp_str = line.split("Optimized speedup:")[1].split("x")[0].strip()
                data[current_kernel]["iterations"].append(float(sp_str))
            elif "FAILED correctness" in line and "Optimized kernel" in line and current_kernel:
                data[current_kernel]["iterations"].append("FAIL")

    return data


def build_slides(data_1x, data_8x, log_1x, log_8x):
    """Build all slides."""

    prof_1x = extract_profiling_data(log_1x)
    prof_8x = extract_profiling_data(log_8x)

    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), prs.slide_width, Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Profiling-Guided Kernel Optimization"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "NCU Profiling Feedback for Iterative Triton Kernel Improvement"
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Xiao Qin  |  University of Leeds  |  April 2026"
    p3.font.size = Pt(18)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 2: Method Overview ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Method: Profiling-Feedback Optimization Loop")

    add_text_box(slide, 0.6, 1.4, 12, 5.5, [
        ("Two-phase pipeline:", True, DARK_BLUE),
        ("", False, None),
        ("Phase 1 -- Correctness (existing, up to 10 attempts)", True, DARK_BLUE),
        ("  C code -> compiler analysis -> structured prompt -> LLM -> Triton kernel", False, None),
        ("  Retry with error feedback until correctness passes", False, GRAY),
        ("", False, None),
        ("Phase 2 -- Profiling Optimization (NEW, up to 3 iterations)", True, GREEN),
        ("  After a correct kernel exists:", False, None),
        ("  1. Profile with NVIDIA Nsight Compute (NCU)", False, None),
        ("     - SM throughput, memory throughput, occupancy, L1/L2 hit rates", False, GRAY),
        ("     - Classify bottleneck: compute-bound / memory-bound / latency-bound", False, GRAY),
        ("  2. Feed NCU metrics + bottleneck diagnosis + current code to LLM", False, None),
        ("  3. LLM generates optimized version", False, None),
        ("  4. Re-verify correctness (reject if broken)", False, None),
        ("  5. Re-benchmark, keep only if speedup improved", False, None),
        ("  6. If improved, re-profile the new version for next iteration", False, None),
        ("     If not improved, reuse cached NCU profile (no redundant profiling)", False, GRAY),
    ], font_size=14)

    # ========== SLIDE 3: 1x Results Table ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Results: 1x Scale (Standard PolyBench Sizes)", "N = 60-120")

    header = ["Kernel", "Pass?", "Attempts", "Initial", "Final", "Change", "Bottleneck", "SM%", "Mem%", "Occ%"]
    rows = [header]
    improved_count = 0
    total_benchmarked = 0

    for kernel in sorted(prof_1x.keys()):
        d = prof_1x[kernel]
        init = d.get("initial_speedup")
        final = d.get("final_speedup")
        ncu = d.get("ncu", {})
        passed = "Y" if init is not None else "N"

        if init is not None:
            total_benchmarked += 1
            init_str = f"{init:.2f}x"
            final_str = f"{final:.2f}x" if final else "-"
            if final and final > init:
                change = f"+{((final/init)-1)*100:.0f}%"
                improved_count += 1
            else:
                change = "-"
        else:
            init_str = "-"
            final_str = "-"
            change = "-"

        bottleneck = ncu.get("bottleneck", "-")
        if "memory" in bottleneck:
            bottleneck = "mem-bound"
        elif "compute" in bottleneck:
            bottleneck = "compute"
        elif "latency" in bottleneck:
            bottleneck = "latency"
        else:
            bottleneck = bottleneck[:12]

        rows.append([
            kernel, passed, str(d.get("initial_speedup") is not None and 1 or "-"),
            init_str, final_str, change,
            bottleneck,
            ncu.get("sm", "-"), ncu.get("mem", "-"), ncu.get("occ", "-"),
        ])

    if len(rows) > 1:
        add_table(slide, rows[:17], 0.3, 1.3, 12.7,
                  col_widths=[1.8, 0.6, 0.8, 1.0, 1.0, 0.9, 1.4, 0.8, 0.8, 0.8])

    # ========== SLIDE 4: 1x Results Table (continued) ==========
    if len(rows) > 17:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_bar(slide, "Results: 1x Scale (continued)")
        remaining = [header] + rows[17:]
        add_table(slide, remaining, 0.3, 1.3, 12.7,
                  col_widths=[1.8, 0.6, 0.8, 1.0, 1.0, 0.9, 1.4, 0.8, 0.8, 0.8])

    # ========== SLIDE 5: 8x Results Table ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Results: 8x Scale (Larger Data Sizes)", "N = 480-960")

    header8 = ["Kernel", "Pass?", "Initial", "Final", "Change", "Bottleneck", "SM%", "Mem%", "Occ%", "L1%"]
    rows8 = [header8]
    improved_8x = 0
    total_8x = 0

    for kernel in sorted(prof_8x.keys()):
        d = prof_8x[kernel]
        init = d.get("initial_speedup")
        final = d.get("final_speedup")
        ncu = d.get("ncu", {})
        passed = "Y" if init is not None else "N"

        if init is not None:
            total_8x += 1
            init_str = f"{init:.1f}x"
            final_str = f"{final:.1f}x" if final else "-"
            if final and final > init:
                change = f"+{((final/init)-1)*100:.0f}%"
                improved_8x += 1
            else:
                change = "-"
        else:
            init_str = "-"
            final_str = "-"
            change = "-"

        bottleneck = ncu.get("bottleneck", "-")
        if "memory" in bottleneck:
            bottleneck = "mem-bound"
        elif "compute" in bottleneck:
            bottleneck = "compute"
        elif "latency" in bottleneck:
            bottleneck = "latency"
        else:
            bottleneck = bottleneck[:12]

        rows8.append([
            kernel, passed, init_str, final_str, change,
            bottleneck,
            ncu.get("sm", "-"), ncu.get("mem", "-"), ncu.get("occ", "-"), ncu.get("l1", "-"),
        ])

    if len(rows8) > 1:
        add_table(slide, rows8[:17], 0.3, 1.3, 12.7,
                  col_widths=[1.8, 0.6, 1.0, 1.0, 0.9, 1.3, 0.8, 0.8, 0.8, 0.8])

    if len(rows8) > 17:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_bar(slide, "Results: 8x Scale (continued)")
        remaining8 = [header8] + rows8[17:]
        add_table(slide, remaining8, 0.3, 1.3, 12.7,
                  col_widths=[1.8, 0.6, 1.0, 1.0, 0.9, 1.3, 0.8, 0.8, 0.8, 0.8])

    # ========== SLIDE: Summary Statistics ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Summary: Profiling-Feedback Impact")

    # Compute stats
    improvements_1x = [(k, prof_1x[k]) for k in prof_1x if prof_1x[k].get("improved")]
    improvements_8x = [(k, prof_8x[k]) for k in prof_8x if prof_8x[k].get("improved")]

    lines_1x = [f"  {k}: {d['initial_speedup']:.2f}x -> {d['final_speedup']:.2f}x (+{((d['final_speedup']/d['initial_speedup'])-1)*100:.0f}%)" for k, d in improvements_1x]
    lines_8x = [f"  {k}: {d['initial_speedup']:.1f}x -> {d['final_speedup']:.1f}x (+{((d['final_speedup']/d['initial_speedup'])-1)*100:.0f}%)" for k, d in improvements_8x]

    add_info_box(slide, 0.5, 1.4, 5.8, 2.8, DARK_BLUE, LIGHT_BLUE_BG,
        f"1x Scale: {len(improvements_1x)} kernels improved", [
        f"Out of {len([k for k in prof_1x if prof_1x[k].get('initial_speedup')])} benchmarked kernels",
        "",
    ] + (lines_1x if lines_1x else ["  (none)"]))

    add_info_box(slide, 6.8, 1.4, 6.0, 2.8, DARK_BLUE, LIGHT_BLUE_BG,
        f"8x Scale: {len(improvements_8x)} kernels improved", [
        f"Out of {len([k for k in prof_8x if prof_8x[k].get('initial_speedup')])} benchmarked kernels",
        "",
    ] + (lines_8x if lines_8x else ["  (none)"]))

    # Failure analysis
    fail_correctness_1x = sum(1 for k in prof_1x if any(i == "FAIL" for i in prof_1x[k].get("iterations", [])))
    fail_correctness_8x = sum(1 for k in prof_8x if any(i == "FAIL" for i in prof_8x[k].get("iterations", [])))

    add_info_box(slide, 0.5, 4.5, 12.3, 2.5, RED, LIGHT_RED_BG,
        "Why most optimizations fail", [
        "",
        f"Kernels where at least one optimization broke correctness: {fail_correctness_1x} (1x), {fail_correctness_8x} (8x)",
        "",
        "The LLM frequently breaks correctness when attempting performance optimizations.",
        "This is the key limitation of single-shot prompting for optimization.",
        "RL-based methods (Dr. Kernel, FT GPT-5) and search (K-Search, OptiML) address this",
        "by exploring many candidates — our correctness check catches failures but can't fix them.",
    ])

    # ========== SLIDE: Key Findings ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Key Findings & Next Steps")

    add_text_box(slide, 0.6, 1.4, 5.5, 5.5, [
        ("Findings", True, DARK_BLUE),
        ("", False, None),
        ("1. Profiling feedback CAN improve performance", True, GREEN),
        ("   NCU metrics help LLM identify and fix bottlenecks", False, None),
        ("", False, None),
        ("2. Correctness is the main bottleneck", True, RED),
        ("   LLM frequently breaks correctness during optimization", False, None),
        ("   Single-shot prompting insufficient for reliable optimization", False, None),
        ("", False, None),
        ("3. Already-optimized kernels resist improvement", True, DARK_BLUE),
        ("   Kernels near HW limits (>80% mem, >90% occ) can't improve", False, None),
        ("", False, None),
        ("4. Small data sizes limit profiling usefulness", True, DARK_BLUE),
        ("   At 1x, most kernels are latency-bound (launch overhead)", False, None),
        ("   At 8x, real bottlenecks emerge and feedback is more useful", False, None),
    ], font_size=14)

    add_info_box(slide, 6.5, 1.4, 6.3, 5.5, GREEN, LIGHT_GREEN_BG,
        "Next Steps for Conference Paper", [
        "",
        "Our compiler analysis + profiling feedback is a foundation.",
        "To overcome the correctness bottleneck, combine with:",
        "",
        "1. Search-based optimization (K-Search / MCTS style)",
        "   - Generate many candidates, keep best correct one",
        "   - Use compiler analysis to verify parallelism safety",
        "",
        "2. Multi-turn RL (Dr. Kernel style)",
        "   - Use NCU metrics as RL reward signal",
        "   - Use correctness check as constraint",
        "",
        "3. Our unique angle: compiler analysis as",
        "   correctness oracle + profiling as performance oracle",
        "   = both signals that RL/search methods need",
    ])

    return prs


if __name__ == "__main__":
    base = Path("/home/qinxiao/workspace/compiler-guided-triton-gen")

    results_1x = base / "results" / "polybench" / "polybench_results" / "results.json"
    results_8x = base / "results" / "polybench" / "polybench_results_scale8x" / "results.json"
    log_1x = Path("/tmp/profiling_1x_full.log")
    log_8x = Path("/tmp/profiling_8x_full.log")

    print(f"1x log: {log_1x} ({'exists' if log_1x.exists() else 'missing'})")
    print(f"8x log: {log_8x} ({'exists' if log_8x.exists() else 'missing'})")

    data_1x = load_results(results_1x) if results_1x.exists() else {}
    data_8x = load_results(results_8x) if results_8x.exists() else {}

    prs = build_slides(data_1x, data_8x, log_1x, log_8x)

    output = base / "presentation" / "profiling_feedback_slides.pptx"
    prs.save(str(output))
    print(f"Saved to {output}")
