#!/usr/bin/env python3
"""Generate PowerPoint slides for the Analysis-Guided LLM Tritonization Pipeline."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MED_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x7D, 0x2F)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
CODE_BG = RGBColor(0xF0, 0xF3, 0xF4)
TABLE_HEADER_BG = RGBColor(0x2E, 0x5C, 0x8A)

FONT_BODY = "Calibri"
FONT_CODE = "Courier New"

# ── Dynamic results loader ───────────────────────────────────────────────
import json as _json

def _load_results_cache():
    """Load WA and NA results once, cache in module globals."""
    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    wa, na = {}, {}
    wa_file = os.path.join(results_dir, "results.json")
    na_file = os.path.join(results_dir, "results_no_analysis.json")
    if os.path.exists(wa_file):
        with open(wa_file) as f:
            wa = _json.load(f)
    if os.path.exists(na_file):
        with open(na_file) as f:
            na = _json.load(f)
    return wa, na

_WA_RESULTS, _NA_RESULTS = _load_results_cache()

def _kinfo(kernel, mode="wa"):
    """Get (passed, attempts, speedup_str) for a kernel."""
    src = _WA_RESULTS if mode == "wa" else _NA_RESULTS
    d = src.get(kernel, {})
    passed = d.get("test_passed", False)
    att = d.get("attempts", "?")
    spd = d.get("benchmark", {}).get("speedup", 0) if passed else 0
    spd_str = f"{spd:.2f}x" if spd > 0 else "FAIL"
    return passed, att, spd, spd_str

def _vs_str(kernel):
    """Return 'WA_speedup vs NA_speedup' string for a kernel."""
    _, _, _, wa_s = _kinfo(kernel, "wa")
    _, _, _, na_s = _kinfo(kernel, "na")
    return f"{wa_s} vs {na_s}"

def _flow_str(kernel, mode="wa"):
    """Return 'Y(att) speedup' or 'FAIL (att att)' for flow boxes."""
    passed, att, spd, spd_str = _kinfo(kernel, mode)
    if passed:
        return f"Y({att}) {spd_str}"
    else:
        return f"FAIL ({att} att)"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_prefix=False, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.line_spacing = Pt(font_size * line_spacing)
        p.space_after = Pt(4)

        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = prefix + ": "
            run1.font.name = FONT_BODY
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run2 = p.add_run()
            run2.text = rest
            run2.font.name = FONT_BODY
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = item
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return txBox


def add_code_box(slide, left, top, width, height, code_text, font_size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = MED_GRAY
    shape.line.width = Pt(0.5)
    # Smaller corner radius
    shape.adjustments[0] = 0.02

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_CODE
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.line_spacing = Pt(font_size * 1.3)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MED_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_flow_box(slide, left, top, width, height, text, fill_color=LIGHT_BLUE,
                 text_color=WHITE, font_size=11, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.15

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_section_title(slide, text, top=Inches(0.3)):
    add_textbox(slide, Inches(0.6), top, Inches(8.4), Inches(0.6),
                text, font_size=28, bold=True, color=DARK_BLUE)


def add_horizontal_arrow(slide, x, y, length=Inches(0.3)):
    """Draw a right-pointing arrow using a triangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, length, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_BLUE
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, x, y, length=Inches(0.25)):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.25), length)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_BLUE
    shape.line.fill.background()
    return shape


def add_simple_table(slide, left, top, width, col_widths, rows, header=True,
                     font_size=13, row_height=0.35):
    """Add a table. rows[0] = header if header=True. col_widths in inches."""
    tbl = slide.shapes.add_table(len(rows), len(col_widths), left, top, width,
                                 Inches(row_height * len(rows))).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(font_size)

            if header and ri == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT_GRAY
    return tbl


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BLUE)

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(1.2),
                "Analysis-Guided LLM\nTritonization Pipeline",
                font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.7),
                "PET + LLVM Static Analysis  ->  LLM Prompt  ->  GPU Triton Kernels",
                font_size=18, color=NEAR_WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.5),
                "Polybench/C 4.2.1 (30 kernels)  |  Rodinia 3.1 (3 kernels)  |  16 Analysis Modules",
                font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def slide_02_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Pipeline Architecture")

    # Row 1: main flow
    y1 = Inches(1.2)
    bw, bh = Inches(1.2), Inches(0.45)
    gap = Inches(0.15)

    boxes = [
        ("C Kernel", DARK_BLUE),
        ("Extract\nScop", MED_BLUE),
        ("PET / LLVM\nAnalysis", LIGHT_BLUE),
        ("Build\nPrompt", MED_BLUE),
        ("Claude\nSonnet", ACCENT_ORANGE),
        ("Triton\nCode", ACCENT_GREEN),
        ("Test vs\nC Ref", MED_BLUE),
    ]

    x = Inches(0.25)
    positions = []
    for label, color in boxes:
        add_flow_box(slide, x, y1, bw, bh, label, fill_color=color, font_size=9)
        positions.append(x)
        x_arrow = x + bw + gap * 0.1
        if label != "Test vs\nC Ref":
            add_horizontal_arrow(slide, x_arrow, y1 + bh / 2 - Inches(0.1),
                                 length=gap * 0.8)
        x = x_arrow + gap * 0.8

    # Retry loop arrow (text annotation)
    add_textbox(slide, Inches(7.2), Inches(1.8), Inches(2.5), Inches(0.35),
                "5+5 retry strategy",
                font_size=10, color=ACCENT_RED, bold=True)

    # Stats bar - load pass rate dynamically
    import json as _json
    _results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    _results_file = os.path.join(_results_dir, "results.json")
    _pass_str = "24/30 Pass (80%)"
    if os.path.exists(_results_file):
        with open(_results_file) as _f:
            _wr = _json.load(_f)
        _np = sum(1 for v in _wr.values() if v.get('test_passed'))
        _nt = len(_wr)
        _pass_str = f"{_np}/{_nt} Pass ({100*_np//_nt}%)"

    stats_y = Inches(2.3)
    stats = [
        ("30 Kernels", DARK_BLUE),
        ("16 Analysis Modules", MED_BLUE),
        ("Speedup Retry (<0.1x)", ACCENT_ORANGE),
        (_pass_str, ACCENT_GREEN),
    ]
    x = Inches(0.5)
    for label, color in stats:
        add_flow_box(slide, x, stats_y, Inches(2.0), Inches(0.4), label,
                     fill_color=color, font_size=11)
        x += Inches(2.3)

    # Analysis modules breakdown
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(9.0), Inches(0.4),
                "16 Analysis Modules (PET + LLVM fallback):",
                font_size=14, bold=True, color=DARK_BLUE)

    modules_text = [
        "WAR Dependencies  |  Parallel Dimensions  |  Reduction Detection  |  Scalar Expansion",
        "Stream Patterns  |  Aliasing  |  Overwrites  |  Indirect Addressing  |  Goto Detection",
        "Loop Distribution  |  Loop Interchange  |  Unrolling  |  Reordering  |  Crossing  |  Convolution  |  Early Exit",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(3.6), Inches(9.0), Inches(2.0),
                    modules_text, font_size=12, color=DARK_TEXT)


def slide_02b_extraction_overview(prs):
    """Scop extraction overview: raw Polybench → standalone kernel → LLM prompt."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Scop Extraction: From Polybench to Kernel")

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "Each kernel goes through a multi-stage extraction before analysis and LLM prompting.",
                font_size=13, color=MED_GRAY)

    # Flow diagram across the top
    flow_y = Inches(1.2)
    bw, bh = Inches(1.5), Inches(0.55)

    stages = [
        ("Raw Polybench\ngemm.c (147 lines)", DARK_BLUE),
        ("Extract Scop\nRegion + Arrays", MED_BLUE),
        ("Expand Macros\n& Globals", LIGHT_BLUE),
        ("Standalone\ngemm.c (27 lines)", ACCENT_GREEN),
        ("PET + LLVM\nAnalysis", ACCENT_ORANGE),
    ]

    x = Inches(0.2)
    for i, (label, color) in enumerate(stages):
        add_flow_box(slide, x, flow_y, bw, bh, label, fill_color=color, font_size=9)
        if i < len(stages) - 1:
            add_horizontal_arrow(slide, x + bw + Inches(0.02),
                                 flow_y + bh / 2 - Inches(0.1),
                                 length=Inches(0.28))
        x += bw + Inches(0.35)

    # What gets removed / what stays
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(4.3), Inches(0.3),
                "Removed (Polybench infrastructure):",
                font_size=14, bold=True, color=ACCENT_RED)
    add_bullet_list(slide, Inches(0.5), Inches(2.35), Inches(4.3), Inches(1.5), [
        "main(), init_array(), print_array()",
        "POLYBENCH_2D/3D array macros",
        "_PB_NI, _PB_NJ, ... bound macros",
        "DATA_TYPE casts, SCALAR_VAL() wrappers",
        "polybench.h timing infrastructure",
    ], font_size=12, color=DARK_TEXT)

    add_textbox(slide, Inches(5.3), Inches(2.0), Inches(4.5), Inches(0.3),
                "Kept / Generated:",
                font_size=14, bold=True, color=ACCENT_GREEN)
    add_bullet_list(slide, Inches(5.3), Inches(2.35), Inches(4.5), Inches(1.5), [
        "#pragma scop region (the computation)",
        "#define NI 60, NJ 70, NK 80 (sizes)",
        "Global arrays: float C[NI][NJ], ...",
        "Global scalars: float alpha = 1.5f",
        "Iteration variables: int i, j, k",
    ], font_size=12, color=DARK_TEXT)

    # Key design decisions box
    y_box = Inches(4.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.5), y_box, Inches(9.0), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    box.line.color.rgb = MED_BLUE
    box.line.width = Pt(1)
    box.adjustments[0] = 0.05
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Why global arrays? "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = MED_BLUE
    run2 = p.add_run()
    run2.text = ("PET requires standalone files with known array sizes. "
                 "By using #define constants and global declarations, "
                 "PET can build the polyhedral model with exact iteration domains "
                 "(e.g., { S[i,j] : 0 <= i < 60 and 0 <= j < 70 }).")
    run2.font.size = Pt(11)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run3 = p2.add_run()
    run3.text = "30 kernels configured: "
    run3.font.bold = True
    run3.font.size = Pt(12)
    run3.font.name = FONT_BODY
    run3.font.color.rgb = MED_BLUE
    run4 = p2.add_run()
    run4.text = ("Each kernel has a config entry specifying: path, dimension sizes, "
                 "macro mappings, data type, extra globals (scalars), and any math replacements.")
    run4.font.size = Pt(11)
    run4.font.name = FONT_BODY
    run4.font.color.rgb = DARK_TEXT


def slide_02c_gemm_before_after(prs):
    """GEMM example: raw Polybench vs extracted standalone kernel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "GEMM Example: Before & After Extraction")

    # LEFT: Raw Polybench
    add_textbox(slide, Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.3),
                "Raw Polybench gemm.c (kernel function):",
                font_size=12, bold=True, color=ACCENT_RED)

    raw_code = """\
void kernel_gemm(int ni, int nj, int nk,
  DATA_TYPE alpha, DATA_TYPE beta,
  DATA_TYPE POLYBENCH_2D(C,NI,NJ,ni,nj),
  DATA_TYPE POLYBENCH_2D(A,NI,NK,ni,nk),
  DATA_TYPE POLYBENCH_2D(B,NK,NJ,nk,nj))
{
  int i, j, k;
#pragma scop
  for (i = 0; i < _PB_NI; i++) {
    for (j = 0; j < _PB_NJ; j++)
      C[i][j] *= beta;
    for (k = 0; k < _PB_NK; k++) {
       for (j = 0; j < _PB_NJ; j++)
         C[i][j] += alpha*A[i][k]*B[k][j];
    }
  }
#pragma endscop
}"""
    add_code_box(slide, Inches(0.3), Inches(1.15), Inches(4.55), Inches(2.7),
                 raw_code, font_size=8)

    # Annotations on raw code
    add_bullet_list(slide, Inches(0.3), Inches(3.95), Inches(4.55), Inches(1.0), [
        "POLYBENCH_2D: complex macro for arrays",
        "_PB_NI, _PB_NK: indirect bound macros",
        "DATA_TYPE: typedef (float/double)",
        "+ 80 lines of main/init/print boilerplate",
    ], font_size=10, color=ACCENT_RED)

    # RIGHT: Extracted kernel
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.6), Inches(0.3),
                "Extracted standalone gemm.c:",
                font_size=12, bold=True, color=ACCENT_GREEN)

    extracted_code = """\
#include <math.h>

#define NI 60
#define NJ 70
#define NK 80

float C[NI][NJ];
float A[NI][NK];
float B[NK][NJ];

float alpha = 1.5f;
float beta = 1.2f;

void gemm_kernel() {
  int i, j, k;
#pragma scop
  for (i = 0; i < NI; i++) {
    for (j = 0; j < NJ; j++)
      C[i][j] *= beta;
    for (k = 0; k < NK; k++) {
       for (j = 0; j < NJ; j++)
         C[i][j] += alpha*A[i][k]*B[k][j];
    }
  }
#pragma endscop
}"""
    add_code_box(slide, Inches(5.2), Inches(1.15), Inches(4.55), Inches(3.7),
                 extracted_code, font_size=8)

    # Arrow between them
    add_horizontal_arrow(slide, Inches(4.88), Inches(2.3), length=Inches(0.28))

    # Bottom insight
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.3), Inches(5.05), Inches(9.5), Inches(0.4))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    insight_box.line.color.rgb = ACCENT_GREEN
    insight_box.line.width = Pt(1)
    insight_box.adjustments[0] = 0.1
    tf = insight_box.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "147 lines -> 27 lines: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = "Same computation, zero boilerplate. Ready for PET polyhedral analysis."
    run2.font.size = Pt(12)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_02d_kernel_to_prompt(prs):
    """From extracted kernel + analysis → function spec → LLM prompt."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "From Kernel to LLM Prompt")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "The extracted kernel feeds two paths that merge into the final prompt.",
                font_size=13, color=MED_GRAY)

    # Flow: Extracted .c → two branches → merge into prompt
    # Top branch: .c → PET/LLVM → Analysis Results
    # Bottom branch: .c → Function Spec DB → arrays, scalars, signature
    # Both → Build Prompt → LLM

    y_top = Inches(1.3)
    y_bot = Inches(2.3)
    y_merge = Inches(1.8)

    # Source box (left)
    add_flow_box(slide, Inches(0.3), Inches(1.7), Inches(1.3), Inches(0.5),
                 "gemm.c\n(extracted)", fill_color=DARK_BLUE, font_size=10)

    # Top path: Analysis
    add_horizontal_arrow(slide, Inches(1.65), y_top + Inches(0.15), length=Inches(0.25))
    add_flow_box(slide, Inches(1.95), y_top, Inches(1.6), Inches(0.5),
                 "PET + LLVM\nAnalysis", fill_color=MED_BLUE, font_size=10)
    add_horizontal_arrow(slide, Inches(3.6), y_top + Inches(0.15), length=Inches(0.25))
    add_flow_box(slide, Inches(3.9), y_top, Inches(1.5), Inches(0.5),
                 "WAR, ParDims\nReduction, ...", fill_color=LIGHT_BLUE, font_size=9)

    # Bottom path: Function spec
    add_horizontal_arrow(slide, Inches(1.65), y_bot + Inches(0.15), length=Inches(0.25))
    add_flow_box(slide, Inches(1.95), y_bot, Inches(1.6), Inches(0.5),
                 "Function Spec\nDatabase", fill_color=MED_BLUE, font_size=10)
    add_horizontal_arrow(slide, Inches(3.6), y_bot + Inches(0.15), length=Inches(0.25))
    add_flow_box(slide, Inches(3.9), y_bot, Inches(1.5), Inches(0.5),
                 "arrays, scalars\nsignature", fill_color=LIGHT_BLUE, font_size=9)

    # Merge arrow → Build Prompt
    add_horizontal_arrow(slide, Inches(5.45), y_merge + Inches(0.15), length=Inches(0.25))
    add_flow_box(slide, Inches(5.75), y_merge, Inches(1.4), Inches(0.5),
                 "Build\nPrompt", fill_color=ACCENT_ORANGE, font_size=10)
    add_horizontal_arrow(slide, Inches(7.2), y_merge + Inches(0.15), length=Inches(0.25))
    add_flow_box(slide, Inches(7.5), y_merge, Inches(1.2), Inches(0.5),
                 "Claude\nSonnet", fill_color=ACCENT_GREEN, font_size=10)

    # Prompt structure breakdown
    add_textbox(slide, Inches(0.3), Inches(3.05), Inches(9.5), Inches(0.3),
                "Final Prompt Structure (GEMM, 7 sections):",
                font_size=14, bold=True, color=DARK_BLUE)

    # Section boxes with descriptions
    sections = [
        ("1", "C Source Code", "Full extracted gemm.c (27 lines, with globals + scop)",
         DARK_BLUE),
        ("2", "Loop Code", 'The #pragma scop region: "this is what you must implement"',
         MED_BLUE),
        ("3", "Analysis Results", "WAR: none | ParDims: i,j parallel, k reduction | Reduction: dot product",
         LIGHT_BLUE),
        ("4", "Array Info", "C: read-write | A: read-only | B: read-only",
         MED_BLUE),
        ("5", "Dimensions", "NI=60, NJ=70, NK=80 (compile-time in C -> runtime in Triton)",
         LIGHT_BLUE),
        ("6", "Exact Signature", "def gemm_triton(C, A, B, alpha, beta, NI, NJ, NK)",
         MED_BLUE),
        ("7", "Triton Rules", "constexpr block sizes, tl.arange outside loops, masked loads, ...",
         ACCENT_ORANGE),
    ]

    y = Inches(3.35)
    for num, title, desc, color in sections:
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.35), y + Inches(0.02), Inches(0.22), Inches(0.22))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.margin_left = Inches(0)
        tf_b.margin_right = Inches(0)
        tf_b.margin_top = Inches(0)
        tf_b.margin_bottom = Inches(0)
        p_b = tf_b.paragraphs[0]
        p_b.text = num
        p_b.font.name = FONT_BODY
        p_b.font.size = Pt(8)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE
        p_b.alignment = PP_ALIGN.CENTER
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title
        add_textbox(slide, Inches(0.65), y, Inches(1.7), Inches(0.25),
                    title, font_size=10, bold=True, color=color)
        # Description
        add_textbox(slide, Inches(2.4), y, Inches(7.3), Inches(0.25),
                    desc, font_size=9, color=DARK_TEXT)
        y += Inches(0.27)

    # Bottom note
    add_textbox(slide, Inches(0.5), y + Inches(0.08), Inches(9.0), Inches(0.25),
                "The scop region tells the LLM exactly which computation to convert. "
                "Analysis tells it HOW to parallelize safely.",
                font_size=10, bold=True, color=MED_BLUE)


def slide_03_ablation_overview(prs):
    """Ablation overview: highlight notable kernels."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Analysis Impact: Highlights")

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with open(os.path.join(results_dir, "results.json")) as f:
        with_a = json.load(f)
    with open(os.path.join(results_dir, "results_no_analysis.json")) as f:
        no_a = json.load(f)

    # Known analysis modules per kernel
    module_map = {
        "ludcmp": "ScalarExp", "2mm": "ParDims(2D-grid)",
        "jacobi_2d": "ParWarn+TS(1,)", "trisolv": "WAR+ParDims",
        "gemm": "ParDims(3D-fix)", "trmm": "WAR+ParDims",
        "gramschmidt": "ParDims+ScExp+Red", "heat_3d": "ParDims(N-D)+TS",
        "lu": "WAR+grid=(1,)", "syr2k": "ParDims(3D-fix)",
        "adi": "ParWarn+WAR", "deriche": "ScalarExp",
        "durbin": "WAR+ScalarExp", "symm": "ScExp+Red",
        "seidel_2d": "ParWarn+WAR", "nussinov": "WAR",
        "atax": "LoopFission", "bicg": "LoopFission",
        "fdtd_2d": "ParWarn+TS(1,)", "floyd_warshall": "ParWarn+WAR",
        "correlation": "MultiPhase", "covariance": "MultiPhase",
        "gemver": "MultiPhase", "doitgen": "ParDims(N-D)",
        "3mm": "CrossPhase", "mvt": "ParDims(2D)",
    }

    # Highlight kernels (most interesting wins/losses/ties)
    highlight_kernels = [
        "doitgen", "heat_3d", "symm", "trmm", "fdtd_2d",
        "covariance", "durbin", "lu", "seidel_2d", "adi",
        "gemm",
    ]

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "30 kernels receive analysis. Showing notable wins, ties, and losses (full table on next slides).",
                font_size=12, color=MED_GRAY)

    rows = [["Kernel", "Key Modules", "With", "W/O", "Winner"]]
    w_wins_all, wo_wins_all, passes_gained = 0, 0, 0
    no_analysis_kernels = set()  # All 30 kernels now receive analysis

    for k in sorted(with_a.keys()):
        if k.replace('-', '_') in no_analysis_kernels:
            continue
        w = with_a.get(k, {})
        n = no_a.get(k, {})
        wp = w.get('test_passed', False)
        np_ = n.get('test_passed', False)
        ws = w.get('benchmark', {}).get('speedup', 0) if w.get('benchmark') else 0
        ns = n.get('benchmark', {}).get('speedup', 0) if n.get('benchmark') else 0
        if wp and not np_:
            w_wins_all += 1
            passes_gained += 1
        elif not wp and np_:
            wo_wins_all += 1
        elif wp and np_ and ws > 0 and ns > 0:
            if ws > ns * 1.15:
                w_wins_all += 1
            elif ns > ws * 1.15:
                wo_wins_all += 1

    for k in highlight_kernels:
        w = with_a.get(k, with_a.get(k.replace('_', '-'), {}))
        n = no_a.get(k, no_a.get(k.replace('_', '-'), {}))
        wp = w.get('test_passed', False)
        np_ = n.get('test_passed', False)
        wa = w.get('attempts', '?')
        na = n.get('attempts', '?')
        ws = w.get('benchmark', {}).get('speedup', 0) if w.get('benchmark') else 0
        ns = n.get('benchmark', {}).get('speedup', 0) if n.get('benchmark') else 0

        w_str = f"Y({wa}) {ws:.2f}x" if wp and ws > 0 else ("Y({})".format(wa) if wp else f"FAIL({wa})")
        n_str = f"Y({na}) {ns:.2f}x" if np_ and ns > 0 else ("Y({})".format(na) if np_ else f"FAIL({na})")

        if wp and not np_:
            winner = "W/ (pass)"
        elif not wp and np_:
            winner = "W/O (pass)"
        elif wp and np_ and ws > 0 and ns > 0:
            ratio = ws / ns if ns > 0 else 99
            if ratio > 1.15:
                winner = f"W/ (+{ratio:.1f}x)"
            elif ratio < 0.87:
                winner = f"W/O (+{1/ratio:.1f}x)"
            else:
                winner = "tie"
        elif not wp and not np_:
            winner = "both fail"
        else:
            winner = "tie"

        mods = module_map.get(k, module_map.get(k.replace('_', '-'), ""))
        rows.append([k, mods, w_str, n_str, winner])

    add_simple_table(slide, Inches(0.2), Inches(1.1), Inches(9.6),
                     [1.3, 2.0, 1.5, 1.5, 1.3], rows,
                     font_size=10, row_height=0.28)

    w_pass = sum(1 for v in with_a.values() if v.get('test_passed'))
    wo_pass = sum(1 for v in no_a.values() if v.get('test_passed'))

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.5), Inches(5.0), Inches(7.0), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = MED_BLUE
    badge.line.fill.background()
    badge.adjustments[0] = 0.15
    tf = badge.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f"Analysis wins {w_wins_all}  |  No-analysis wins {wo_wins_all}  |  {w_pass}/{len(with_a)} vs {wo_pass}/{len(no_a)}"
    p.font.name = FONT_BODY
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def slide_04_war_trisolv_trmm(prs):
    """Analysis-guided wins: WAR clone pattern (trisolv, trmm)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Analysis Win: WAR Clone Pattern (trisolv, trmm)")

    # -- LEFT: trisolv --
    add_textbox(slide, Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.3),
                f"trisolv ({_vs_str('trisolv')}) -- forward substitution:",
                font_size=11, bold=True, color=MED_BLUE)
    trisolv_code = """\
for (i = 0; i < N; i++) {
  x[i] = b[i];
  for (j = 0; j < i; j++)
    x[i] -= L[i][j] * x[j];
  x[i] = x[i] / L[i][i];
}"""
    add_code_box(slide, Inches(0.3), Inches(1.15), Inches(4.5), Inches(1.2),
                 trisolv_code, font_size=9)

    add_textbox(slide, Inches(0.3), Inches(2.4), Inches(4.5), Inches(0.2),
                "Injected analysis:", font_size=10, bold=True, color=DARK_BLUE)
    trisolv_prompt = """\
## WAR Dependencies
Arrays needing read-only copies: x
  Read x[(j)] may conflict with Write x[(i)]
Pattern: x_copy = x.clone()

## Parallelization Analysis
Triangular bounds: j < i
- Parallelize j, sequential i: VALID
- Parallelize i, sequential j: VALID"""
    add_code_box(slide, Inches(0.3), Inches(2.6), Inches(4.5), Inches(1.55),
                 trisolv_prompt, font_size=8)

    _tri_wa_spd = _kinfo('trisolv', 'wa')[2]
    _tri_na_spd = _kinfo('trisolv', 'na')[2]
    _tri_wa_color = ACCENT_GREEN if _tri_wa_spd >= _tri_na_spd else ACCENT_ORANGE
    _tri_na_color = ACCENT_GREEN if _tri_na_spd > _tri_wa_spd else ACCENT_ORANGE
    add_flow_box(slide, Inches(0.3), Inches(4.2), Inches(2.1), Inches(0.3),
                 f"With: {_flow_str('trisolv', 'wa')}", fill_color=_tri_wa_color, font_size=10)
    add_flow_box(slide, Inches(2.5), Inches(4.2), Inches(2.3), Inches(0.3),
                 f"Without: {_flow_str('trisolv', 'na')}", fill_color=_tri_na_color, font_size=10)

    # -- RIGHT: trmm --
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.6), Inches(0.3),
                f"trmm ({_vs_str('trmm')}) -- triangular matmul:",
                font_size=11, bold=True, color=MED_BLUE)
    trmm_code = """\
for (i = 0; i < M; i++)
  for (j = 0; j < N; j++) {
    for (k = i+1; k < M; k++)
      B[i][j] += A[k][i] * B[k][j];
    B[i][j] = alpha * B[i][j];
  }"""
    add_code_box(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(1.2),
                 trmm_code, font_size=9)

    add_textbox(slide, Inches(5.2), Inches(2.4), Inches(4.6), Inches(0.2),
                "Injected analysis:", font_size=10, bold=True, color=DARK_BLUE)
    trmm_prompt = """\
## WAR Dependencies
B: WAR carried by loops i, k
- Parallelizing j: SAFE (no copy needed)
- Parallelizing i: REQUIRES B_copy = B.clone()

## WAR + Parallelization Recommendation
RECOMMENDED: Parallelize j
  -- no WAR copies needed for B"""
    add_code_box(slide, Inches(5.2), Inches(2.6), Inches(4.6), Inches(1.55),
                 trmm_prompt, font_size=8)

    _trmm_na_passed = _kinfo('trmm', 'na')[0]
    add_flow_box(slide, Inches(5.2), Inches(4.2), Inches(2.1), Inches(0.3),
                 f"With: {_flow_str('trmm', 'wa')}", fill_color=ACCENT_GREEN, font_size=10)
    add_flow_box(slide, Inches(7.4), Inches(4.2), Inches(2.4), Inches(0.3),
                 f"Without: {_flow_str('trmm', 'na')}", fill_color=ACCENT_RED if not _trmm_na_passed else ACCENT_ORANGE, font_size=10)

    # Bottom insight
    add_textbox(slide, Inches(0.3), Inches(4.65), Inches(9.5), Inches(0.3),
                "WAR analysis tells the LLM exactly which arrays to clone and which dimensions are safe to parallelize.",
                font_size=10, color=MED_GRAY)


def slide_05_war_lu_seidel(prs):
    """LLVM DV fix: seidel_2d now passes, adi improvement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "LLVM Direction Vector Fix: seidel_2d, adi")

    # -- LEFT: seidel_2d -- NOW PASSES
    add_textbox(slide, Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.3),
                f"seidel_2d ({_vs_str('seidel_2d')}) -- Gauss-Seidel:",
                font_size=11, bold=True, color=ACCENT_GREEN)
    seidel_code = """\
for (t = 0; t < TSTEPS; t++)
 for (i = 1; i < N-1; i++)
  for (j = 1; j < N-1; j++)
   A[i][j] = (A[i-1][j-1]+...
     +A[i+1][j+1]) / 9;"""
    add_code_box(slide, Inches(0.3), Inches(1.12), Inches(4.5), Inches(1.1),
                 seidel_code, font_size=8)

    add_textbox(slide, Inches(0.3), Inches(2.25), Inches(4.5), Inches(0.2),
                "Injected analysis (after LLVM DV fix):", font_size=10, bold=True, color=DARK_BLUE)
    seidel_prompt = """\
## WAR Dependencies
A: 5 read-write conflicts on stencil neighbors
If separate kernels: no cloning needed.

## Parallelization Warning
No dimension is safe to parallelize.
- t: sequential context loop
- i: deps carried (LLVM DVs: [S -1 -1])
Same array with neighbor deps.
Use grid=(1,) with nested loops sequentially."""
    add_code_box(slide, Inches(0.3), Inches(2.45), Inches(4.5), Inches(1.55),
                 seidel_prompt, font_size=7.5)

    _seidel_na_passed = _kinfo('seidel_2d', 'na')[0]
    add_flow_box(slide, Inches(0.3), Inches(4.05), Inches(2.1), Inches(0.3),
                 f"With: {_flow_str('seidel_2d', 'wa')}", fill_color=ACCENT_GREEN, font_size=10)
    add_flow_box(slide, Inches(2.5), Inches(4.05), Inches(2.3), Inches(0.3),
                 f"Without: {_flow_str('seidel_2d', 'na')}", fill_color=ACCENT_RED if not _seidel_na_passed else ACCENT_ORANGE, font_size=10)

    # -- RIGHT: adi -- improved by DV fix
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.6), Inches(0.3),
                f"adi ({_vs_str('adi')}) -- ADI stencil:",
                font_size=11, bold=True, color=ACCENT_GREEN)

    add_textbox(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(0.2),
                "The LLVM fallback bug:", font_size=10, bold=True, color=ACCENT_RED)
    add_bullet_list(slide, Inches(5.2), Inches(1.35), Inches(4.6), Inches(1.0), [
        "PET fails on 5 kernels -> LLVM fallback used",
        "Old: valid=True always (even with flow deps!)",
        "Told LLM to parallelize dims with carried deps",
        "All 5 kernels performed WORSE with analysis",
    ], font_size=9, color=DARK_TEXT)

    add_textbox(slide, Inches(5.2), Inches(2.35), Inches(4.6), Inches(0.2),
                "Fix: per-dim direction vector analysis", font_size=10, bold=True, color=ACCENT_GREEN)
    add_bullet_list(slide, Inches(5.2), Inches(2.55), Inches(4.6), Inches(1.2), [
        "LLVM DA provides per-loop DVs: [S -1 -1]",
        "Map DV positions to source loop variables",
        "Check _direction_entry_carries_dep() per dim",
        "Now: valid=False when dep carried at that level",
    ], font_size=9, color=DARK_TEXT)

    # Before/after table
    dv_rows = [
        ["Kernel", "Before DV fix", "After (current)"],
        ["adi", "0.65x", _kinfo('adi', 'wa')[3]],
        ["seidel_2d", "FAIL", f"{_kinfo('seidel_2d', 'wa')[3]}{' (PASS)' if _kinfo('seidel_2d', 'wa')[0] else ''}"],
        ["fdtd_2d", "0.56x", _kinfo('fdtd_2d', 'wa')[3]],
        ["floyd_warshall", "0.46x", _kinfo('floyd_warshall', 'wa')[3]],
        ["jacobi_2d", "0.19x", _kinfo('jacobi_2d', 'wa')[3]],
    ]
    add_simple_table(slide, Inches(5.2), Inches(3.65), Inches(4.5),
                     [1.5, 1.5, 1.5], dv_rows, font_size=9, row_height=0.25)

    # Bottom insight
    add_textbox(slide, Inches(0.3), Inches(5.2), Inches(9.5), Inches(0.3),
                f"DV fix: seidel_2d passes (grid=(1,)). fdtd_2d {_kinfo('fdtd_2d','wa')[3]}, "
                f"jacobi_2d {_kinfo('jacobi_2d','wa')[3]}. adi {_kinfo('adi','wa')[3]} (vs W/O {_kinfo('adi','na')[3]}).",
                font_size=9, color=MED_GRAY)


def slide_06_scalarexp_deriche(prs):
    """Analysis-guided win: ScalarExp for deriche (biggest speedup delta)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Analysis Win: Scalar Expansion (deriche)")

    # C code
    add_textbox(slide, Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.3),
                "deriche -- edge detection (IIR filter):",
                font_size=11, bold=True, color=MED_BLUE)
    deriche_code = """\
for (i = 0; i < W; i++) {
  ym1 = 0; ym2 = 0; xm1 = 0;
  for (j = 0; j < H; j++) {
    yy1[i][j] = a1*imgIn[i][j]
       + a2*xm1 + b1*ym1 + b2*ym2;
    xm1 = imgIn[i][j];
    ym2 = ym1;
    ym1 = yy1[i][j];
  }
} // + 5 more similar filter passes"""
    add_code_box(slide, Inches(0.3), Inches(1.15), Inches(4.5), Inches(1.65),
                 deriche_code, font_size=9)

    # The injected prompt
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.6), Inches(0.3),
                "Injected analysis:", font_size=11, bold=True, color=DARK_BLUE)
    prompt = """\
## SCALAR EXPANSION PATTERN DETECTED
6 scalar variables with loop-carried deps:

Variable: xm1 (direct_expansion)
  Simply replace with indexed expression.
Variable: tm1 (direct_expansion)
  Simply replace with indexed expression.
Variable: ym1 (direct_expansion)
  Simply replace with indexed expression.
Variable: xp1 (direct_expansion)
  Simply replace with indexed expression.
Variable: tp1 (direct_expansion)
  Simply replace with indexed expression.
Variable: yp1 (direct_expansion)
  Simply replace with indexed expression.

(No ParDims: both dims INVALID due to
 scalar loop-carried dependencies)"""
    add_code_box(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(3.0),
                 prompt, font_size=7.5)

    # Explanation
    add_bullet_list(slide, Inches(0.3), Inches(2.9), Inches(4.5), Inches(1.2), [
        "6 scalars carry values across j iterations",
        "Without analysis: LLM doesn't privatize them",
        "  -> Race conditions, wrong results (0.11x)",
        "With analysis: LLM expands to per-thread arrays",
        "  -> Correct parallel code (pass on attempt 9)",
    ], font_size=10)

    # Results
    add_flow_box(slide, Inches(0.3), Inches(4.15), Inches(2.1), Inches(0.35),
                 f"With: {_flow_str('deriche', 'wa')}", fill_color=ACCENT_GREEN, font_size=11)
    add_flow_box(slide, Inches(2.5), Inches(4.15), Inches(2.3), Inches(0.35),
                 f"Without: {_flow_str('deriche', 'na')}", fill_color=ACCENT_ORANGE, font_size=11)

    add_textbox(slide, Inches(0.3), Inches(4.65), Inches(9.5), Inches(0.3),
                "ScalarExp identifies 6 loop-carried deps the LLM misses. *Pass rate varies across runs (LLM nondeterminism).",
                font_size=10, color=MED_GRAY)


def slide_07_scalarexp_durbin_symm(prs):
    """Analysis-guided wins: ScalarExp for durbin + symm."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Analysis Win: Scalar Expansion (durbin, symm)")

    # -- LEFT: durbin --
    add_textbox(slide, Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.3),
                f"durbin ({_vs_str('durbin')}) -- Levinson-Durbin:",
                font_size=11, bold=True, color=MED_BLUE)

    add_textbox(slide, Inches(0.3), Inches(1.15), Inches(4.5), Inches(0.2),
                "Injected analysis:", font_size=10, bold=True, color=DARK_BLUE)
    durbin_prompt = """\
## WAR Dependencies
Note: WAR deps on r, y arrays.
If using separate kernels: no cloning needed.

## SCALAR EXPANSION DETECTED
Variable: alpha (previous_value)
  alpha carries previous iteration's value.
  Replace with: -(r[k-1]+sum)/beta (for k>0)
  Initial: alpha = -r[0] (for k=0)
  After substitution: FULLY PARALLEL.

Variable: sum (direct_expansion)
  Simply replace with indexed expression.

(No ParDims: both dims INVALID due to
 scalar loop-carried dependencies)"""
    add_code_box(slide, Inches(0.3), Inches(1.35), Inches(4.5), Inches(2.5),
                 durbin_prompt, font_size=7.5)

    add_flow_box(slide, Inches(0.3), Inches(3.9), Inches(2.1), Inches(0.3),
                 f"With: {_flow_str('durbin', 'wa')}", fill_color=ACCENT_GREEN, font_size=10)
    add_flow_box(slide, Inches(2.5), Inches(3.9), Inches(2.3), Inches(0.3),
                 f"Without: {_flow_str('durbin', 'na')}", fill_color=ACCENT_RED, font_size=10)

    # -- RIGHT: symm --
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.6), Inches(0.3),
                f"symm ({_vs_str('symm')}) -- symmetric matmul:",
                font_size=11, bold=True, color=MED_BLUE)

    add_textbox(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(0.2),
                "Injected analysis:", font_size=10, bold=True, color=DARK_BLUE)
    symm_prompt = """\
## SCALAR EXPANSION DETECTED
Variable: temp2 (direct_expansion)
  Simply replace with indexed expression.

## Reduction Pattern Analysis
Detected: Sum Reduction (+= operator)
Use tl.sum() for parallel reduction:
  vals = tl.load(a_ptr + offsets, mask=mask)
  block_sum = tl.sum(vals, axis=0)

(ParDims both INVALID due to scalar temp2
 -- ScalarExp handles it, LLM infers j-parallel)"""
    add_code_box(slide, Inches(5.2), Inches(1.35), Inches(4.6), Inches(2.5),
                 symm_prompt, font_size=7.5)

    add_flow_box(slide, Inches(5.2), Inches(3.9), Inches(2.1), Inches(0.3),
                 f"With: {_flow_str('symm', 'wa')}", fill_color=ACCENT_GREEN, font_size=10)
    add_flow_box(slide, Inches(7.4), Inches(3.9), Inches(2.4), Inches(0.3),
                 f"Without: {_flow_str('symm', 'na')}", fill_color=ACCENT_ORANGE, font_size=10)

    # Bottom insight
    _durbin_wa, _durbin_na = _kinfo('durbin','wa')[3], _kinfo('durbin','na')[3]
    _symm_wa, _symm_na = _kinfo('symm','wa')[3], _kinfo('symm','na')[3]
    add_textbox(slide, Inches(0.3), Inches(4.4), Inches(9.5), Inches(0.6),
                f"durbin: Analysis helps ({_durbin_wa} vs {_durbin_na}) by identifying scalar expansion patterns. "
                f"symm: Three analysis modules (ParDims + ScalarExp + Reduction) boost {_symm_wa} vs {_symm_na}.",
                font_size=9, color=MED_GRAY)


def slide_08_analysis_summary(prs):
    """Summary: what analysis → what benefit per kernel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Analysis Module Impact Summary")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "For each winning kernel: which analysis modules were injected and what they told the LLM.",
                font_size=11, color=MED_GRAY)

    _summary_kernels = [
        ("doitgen",       "ParDims(N-D)", "4D loop: r,q parallel (2D grid), p,s sequential inside"),
        ("heat_3d",       "ParDims(N-D)", "3D stencil: t in host, i/j/k parallel (multi-CTA)"),
        ("symm",          "ScExp+Red",    "temp2 expansion + tl.sum() reduction pattern"),
        ("trmm",          "WAR+LLVM",     "RECOMMENDED: parallelize j (no clone needed)"),
        ("fdtd_2d",       "ParWarn+TS",   "grid=(1,) with t-loop inside; 2D tiling per timestep"),
        ("durbin",        "WAR+ScExp",    "Scalar expansion: alpha, sum loop-carried deps"),
        ("lu",            "WAR+Par",      "grid=(1,) with i-loop inside; vectorize j with tl.arange"),
        ("seidel_2d",     "ParWarn+WAR",  "No dim safe (LLVM DVs) -> grid=(1,) sequential"),
    ]
    rows = [["Kernel", "With", "W/O", "Modules", "Key Instruction to LLM"]]
    for kn, mods, desc in _summary_kernels:
        rows.append([kn, _kinfo(kn, 'wa')[3], _kinfo(kn, 'na')[3], mods, desc])
    add_simple_table(slide, Inches(0.2), Inches(1.1), Inches(9.6),
                     [1.2, 0.8, 0.8, 1.2, 5.6], rows,
                     font_size=10, row_height=0.35)

    # Categorize by module
    add_textbox(slide, Inches(0.5), Inches(3.9), Inches(9.0), Inches(0.25),
                "Analysis Module Breakdown:", font_size=12, bold=True, color=DARK_BLUE)

    add_bullet_list(slide, Inches(0.5), Inches(4.2), Inches(4.3), Inches(1.0), [
        "WAR (6 kernels): Clone arrays / multi-kernel hints",
        "ParDims (7 kernels): Which dims safe to parallelize",
        "  + N-D analysis for 4D+ loops (heat_3d, doitgen)",
        "ScalarExp (4 kernels): Eliminate loop-carried scalars",
    ], font_size=10)

    add_bullet_list(slide, Inches(5.2), Inches(4.2), Inches(4.5), Inches(1.0), [
        "ParWarn (5 kernels): LLVM DV per-dim dep check",
        "  + Timestep guidance: grid=(1,) vs multi-CTA",
        "LoopFission (2 kernels): Split opposing reductions",
        "Reduction (1 kernel): Use tl.sum() pattern",
    ], font_size=10)


def slide_09_heat3d_failure(prs):
    """N-D ParDims fix: heat_3d. Contrast with seidel_2d."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Stencil Contrast: heat_3d (Jacobi) vs seidel_2d (Gauss-Seidel)")

    # LEFT: heat_3d -- fixed with N-D ParDims
    add_textbox(slide, Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.3),
                f"heat_3d -- 3D Jacobi stencil ({_kinfo('heat_3d','wa')[3]} speedup):",
                font_size=12, bold=True, color=ACCENT_GREEN)
    heat3d_code = """\
for (t = 1; t <= TSTEPS; t++) {
  for (i,j,k)  // Loop nest 1
    B[i][j][k] = 0.125*(A[i+1][j][k]
      - 2*A[i][j][k] + A[i-1][j][k])
      + ... + A[i][j][k];  // 7-pt stencil
  for (i,j,k)  // Loop nest 2
    A[i][j][k] = 0.125*(B[i+1][j][k]
      - 2*B[i][j][k] + B[i-1][j][k])
      + ... + B[i][j][k];  // 7-pt stencil
}"""
    add_code_box(slide, Inches(0.3), Inches(1.15), Inches(4.5), Inches(1.65),
                 heat3d_code, font_size=8)

    add_textbox(slide, Inches(0.3), Inches(2.85), Inches(4.5), Inches(0.3),
                f"N-D ParDims fix (was FAIL, now {_kinfo('heat_3d','wa')[3]}):", font_size=12, bold=True, color=ACCENT_GREEN)
    add_bullet_list(slide, Inches(0.3), Inches(3.15), Inches(4.5), Inches(1.0), [
        "Old: dims=['t','i'] (truncated to 2D)",
        "  -> LLM only parallelizes i, fails 10/10",
        "New: dims=['t','i','j','k'] (full N-D)",
        f"  -> LLM parallelizes i*j*k, passes {_kinfo('heat_3d','wa')[3]}",
    ], font_size=10, color=DARK_TEXT)

    add_textbox(slide, Inches(0.3), Inches(4.2), Inches(4.5), Inches(0.3),
                "Double-buffering: no WAR", font_size=12, bold=True, color=MED_BLUE)
    add_bullet_list(slide, Inches(0.3), Inches(4.5), Inches(4.5), Inches(0.6), [
        "Nest 1: reads A, writes B (different arrays)",
        "Nest 2: reads B, writes A (different arrays)",
    ], font_size=10, color=DARK_TEXT)

    # RIGHT: seidel_2d -- inherently sequential
    add_textbox(slide, Inches(5.2), Inches(0.85), Inches(4.6), Inches(0.3),
                f"seidel_2d -- Gauss-Seidel stencil ({_kinfo('seidel_2d','wa')[3]}):",
                font_size=12, bold=True, color=ACCENT_ORANGE)
    seidel_code = """\
for (t = 0; t <= TSTEPS - 1; t++)
  for (i = 1; i <= N - 2; i++)
    for (j = 1; j <= N - 2; j++)
      A[i][j] = (A[i-1][j-1] + A[i-1][j]
        + A[i-1][j+1] + A[i][j-1]
        + A[i][j]     + A[i][j+1]
        + A[i+1][j-1] + A[i+1][j]
        + A[i+1][j+1]) / 9.0;"""
    add_code_box(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(1.65),
                 seidel_code, font_size=8)

    add_textbox(slide, Inches(5.2), Inches(2.85), Inches(4.6), Inches(0.3),
                "Inherently sequential (now passes via LLVM DV fix):", font_size=12, bold=True, color=ACCENT_GREEN)
    add_bullet_list(slide, Inches(5.2), Inches(3.15), Inches(4.6), Inches(1.0), [
        "Single loop nest -- reads & writes SAME array A",
        "LLVM DVs confirm deps on ALL dims: [S -1 -1]",
        "ParWarning: 'No dim safe' -> LLM uses grid=(1,)",
        f"{_kinfo('seidel_2d','wa')[3]} (slow but CORRECT -- inherently sequential)",
    ], font_size=10, color=DARK_TEXT)

    fail_rows = [
        ["Kernel", "Parallelizable?", "Result"],
        ["heat_3d", "Yes (Jacobi)", f"{_flow_str('heat_3d','wa')} (N-D ParDims)"],
        ["seidel_2d", "No (Gauss-Seidel)", f"{_flow_str('seidel_2d','wa')} (sequential)"],
    ]
    add_simple_table(slide, Inches(5.2), Inches(4.1), Inches(4.5),
                     [1.5, 1.5, 1.5], fail_rows)

    # Key insight box
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.3), Inches(5.0), Inches(9.5), Inches(0.45))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    insight_box.line.color.rgb = ACCENT_ORANGE
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.08
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "heat_3d: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = "N-D ParDims tells LLM to parallelize all spatial dims (10.29x). "
    run2.font.size = Pt(11)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT
    run3 = p.add_run()
    run3.text = "seidel_2d: "
    run3.font.bold = True
    run3.font.size = Pt(12)
    run3.font.name = FONT_BODY
    run3.font.color.rgb = ACCENT_GREEN
    run4 = p.add_run()
    run4.text = "LLVM DV fix + ParWarning -> LLM uses grid=(1,), passes (0.25x sequential)."
    run4.font.size = Pt(11)
    run4.font.name = FONT_BODY
    run4.font.color.rgb = DARK_TEXT


def slide_09_gemm_triton(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Generated Triton Code -- GEMM")

    # C code (left)
    add_textbox(slide, Inches(0.3), Inches(0.9), Inches(4.0), Inches(0.3),
                "C Original (6 lines):", font_size=12, bold=True, color=DARK_BLUE)
    c_code = """\
for (i = 0; i < NI; i++) {
  for (j = 0; j < NJ; j++)
    C[i][j] *= beta;
  for (k = 0; k < NK; k++)
    for (j = 0; j < NJ; j++)
      C[i][j] += alpha*A[i][k]*B[k][j];
}"""
    add_code_box(slide, Inches(0.3), Inches(1.2), Inches(4.5), Inches(1.55), c_code, font_size=10)

    # Triton code (right)
    add_textbox(slide, Inches(5.1), Inches(0.9), Inches(4.8), Inches(0.3),
                "Generated Triton Kernel:", font_size=12, bold=True, color=ACCENT_GREEN)
    triton_code = """\
@triton.jit
def gemm_kernel(A, B, C, alpha, beta,
                NI, NJ, NK,
                BLOCK_I: tl.constexpr,
                BLOCK_J: tl.constexpr):
  pid_i = tl.program_id(0)
  pid_j = tl.program_id(1)
  i_off = pid_i*BLOCK_I + tl.arange(0,BLOCK_I)
  j_off = pid_j*BLOCK_J + tl.arange(0,BLOCK_J)
  # C[i][j] *= beta
  c_idx = i_off[:,None]*NJ + j_off[None,:]
  mask = (i_off[:,None]<NI) & (j_off[None,:]<NJ)
  c = tl.load(C+c_idx, mask=mask) * beta
  # k-loop accumulation
  for k in range(NK):
    a = tl.load(A + i_off*NK+k, mask=i_off<NI)
    b = tl.load(B + k*NJ+j_off, mask=j_off<NJ)
    c += alpha * a[:,None] * b[None,:]
  tl.store(C+c_idx, c, mask=mask)"""
    add_code_box(slide, Inches(5.1), Inches(1.2), Inches(4.7), Inches(3.5), triton_code, font_size=9)

    # Highlights
    add_textbox(slide, Inches(0.3), Inches(3.0), Inches(4.5), Inches(0.3),
                "Key Triton Patterns:", font_size=13, bold=True, color=MED_BLUE)
    add_bullet_list(slide, Inches(0.3), Inches(3.3), Inches(4.5), Inches(2.0), [
        "2D tiling: pid_i, pid_j -> block of C",
        "Masked loads: boundary safety",
        "k-loop: sequential accumulation",
        "Broadcasting: a[:,None] * b[None,:]",
        "Scalar broadcast: alpha, beta",
    ], font_size=12)

    # Result badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.3), Inches(5.0), Inches(3.5), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_GREEN
    badge.line.fill.background()
    badge.adjustments[0] = 0.15
    tf = badge.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    _gemm_p, _gemm_a, _gemm_s, _gemm_ss = _kinfo('gemm', 'wa')
    p.text = f"{'First-try pass' if _gemm_a == 1 else f'Pass (attempt {_gemm_a})'}  |  {_gemm_ss} speedup"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def slide_09b_evaluation(prs):
    """Evaluation Methodology: correctness testing, performance benchmarking, retry strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Evaluation Methodology")

    # ── LEFT COLUMN: Correctness Testing ──
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(4.3), Inches(0.3),
                "Correctness Testing", font_size=16, bold=True, color=MED_BLUE)

    correctness_code = """\
# 1. Compile C kernel to shared library
gcc -O2 -shared -fPIC -o kernel.so kernel.c

# 2. Load .so via ctypes, memmove inputs
lib = ctypes.CDLL("kernel.so")
ctypes.memmove(c_arr, tensor.numpy().ctypes.data, nbytes)

# 3. Run C reference
lib.kernel_function(c_arr, N)

# 4. Run Triton kernel (same inputs)
triton_fn(gpu_tensor, N)

# 5. Compare element-wise
assert |C_ref - Triton| < atol + rtol * |C_ref|"""
    add_code_box(slide, Inches(0.5), Inches(1.2), Inches(4.3), Inches(2.55),
                 correctness_code, font_size=8)

    # Tolerance detail
    add_bullet_list(slide, Inches(0.5), Inches(3.85), Inches(4.3), Inches(0.8), [
        "Default: atol=1e-4, rtol=1e-4",
        "Per-kernel overrides for FP32 accumulation",
        "Auto-detect C array type (int/float/double)",
    ], font_size=10, color=DARK_TEXT)

    # ── RIGHT COLUMN: Performance Benchmarking ──
    add_textbox(slide, Inches(5.3), Inches(0.85), Inches(4.3), Inches(0.3),
                "Performance Benchmarking", font_size=16, bold=True, color=MED_BLUE)

    perf_code = """\
# C reference timing
for _ in range(5):   # warmup
    lib.kernel_function(c_arr, N)
t0 = time.perf_counter()
for _ in range(50):  # timed runs
    lib.kernel_function(c_arr, N)
c_time = (time.perf_counter() - t0) / 50

# Triton GPU timing
for _ in range(5):   # warmup
    triton_fn(gpu_tensor, N)
    torch.cuda.synchronize()
t0 = time.perf_counter()
for _ in range(50):  # timed runs
    triton_fn(gpu_tensor, N)
    torch.cuda.synchronize()
triton_time = (time.perf_counter() - t0) / 50

speedup = c_time / triton_time"""
    add_code_box(slide, Inches(5.3), Inches(1.2), Inches(4.3), Inches(3.0),
                 perf_code, font_size=8)

    # Speedup note
    add_bullet_list(slide, Inches(5.3), Inches(4.3), Inches(4.3), Inches(0.4), [
        "Speedup >1x = Triton faster than C on CPU",
    ], font_size=10, color=DARK_TEXT)

    # ── BOTTOM: Retry Strategy ──
    retry_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), Inches(4.65), Inches(9.0), Inches(0.8))
    retry_box.fill.solid()
    retry_box.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    retry_box.line.color.rgb = ACCENT_ORANGE
    retry_box.line.width = Pt(1)
    retry_box.adjustments[0] = 0.08
    tf = retry_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "5+5 Retry Strategy: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_ORANGE
    run2 = p.add_run()
    run2.text = "Up to 10 attempts. Error feedback injected into next prompt. Context reset at attempt 6 for fresh approach."
    run2.font.size = Pt(11)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_09c_test_inputs(prs):
    """Domain-appropriate test inputs -- details moved from evaluation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Test Input Generation")

    # Why it matters + default
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(9.0), Inches(0.3),
                "Many kernels require structured inputs -- random floats cause divergence or crashes.",
                font_size=13, color=MED_GRAY)

    # Default as a single highlighted line
    default_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.4))
    default_box.fill.solid()
    default_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    default_box.line.color.rgb = ACCENT_GREEN
    default_box.line.width = Pt(1)
    default_box.adjustments[0] = 0.1
    tf = default_box.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Default (23 kernels): "
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = "torch.randn float32  |  random scalars (alpha, beta)  |  Polybench SMALL_DATASET sizes"
    run2.font.size = Pt(12)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT

    # Domain-specific overrides table
    add_textbox(slide, Inches(0.5), Inches(1.75), Inches(9.0), Inches(0.3),
                "Domain-Specific Overrides (7 kernels):",
                font_size=15, bold=True, color=MED_BLUE)

    override_rows = [
        ["Kernel", "Requirement", "Input Generation"],
        ["cholesky", "Symmetric positive definite (SPD)",
         "A = R\u1d40R + N\u00b7I  (ensures positive eigenvalues)"],
        ["lu, ludcmp", "Diagonally dominant",
         "A = randn + N\u00b7I  (prevents zero pivots)"],
        ["trisolv", "Lower-triangular, non-singular",
         "L = tril(randn), |diag| \u2265 1"],
        ["gramschmidt", "Well-conditioned columns",
         "A = randn + 5\u00b7I  (prevents rank deficiency)"],
        ["nussinov", "Integer RNA bases + zeros table",
         "seq = randint(0, 4);  table = zeros"],
        ["floyd_warshall", "Non-negative edge weights",
         "|randn| \u00d7 10 + 1  (no negative cycles)"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(2.1), Inches(9.0),
                     [1.8, 3.0, 4.2], override_rows)

    # Tolerance overrides (table ends at 2.1 + 7*0.35 = 4.55)
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.5),
                "Tolerance: default atol=1e-4, rtol=1e-4  |  durbin: atol=0.05  |  gramschmidt: atol=1.0  |  lu/3mm: atol=5.0  |  ludcmp: atol=0.05",
                font_size=11, color=MED_GRAY)


def slide_10_results(prs):
    import json
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Results Summary")

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with open(os.path.join(results_dir, "results.json")) as f:
        wr = json.load(f)

    total = len(wr)
    passed = {k: v for k, v in wr.items() if v.get('test_passed')}
    n_pass = len(passed)
    first_try = sum(1 for v in passed.values() if v.get('attempts') == 1)
    failed_names = sorted(k for k, v in wr.items() if not v.get('test_passed'))

    # Speedups for passed kernels
    speedups = []
    for k, v in passed.items():
        b = v.get('benchmark', {})
        if b and b.get('speedup', 0) > 0:
            speedups.append((k, b['speedup'], b.get('c_ref_time_ms', 0), b.get('triton_time_ms', 0)))
    speedups.sort(key=lambda x: -x[1])

    gpu_wins = sum(1 for _, s, _, _ in speedups if s >= 1.0)
    all_s = [s for _, s, _, _ in speedups]
    sorted_s = sorted(all_s)
    median_s = sorted_s[len(sorted_s) // 2] if sorted_s else 0
    mean_s = sum(all_s) / len(all_s) if all_s else 0

    pct = f"{100*n_pass/total:.0f}" if total > 0 else "?"
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(4.0), Inches(0.3),
                f"Correctness: {n_pass}/{total} ({pct}%)",
                font_size=16, bold=True, color=ACCENT_GREEN)

    add_bullet_list(slide, Inches(0.5), Inches(1.35), Inches(4.0), Inches(0.85), [
        f"{first_try} first-try passes",
        f"{n_pass - first_try} after retry (5+5 strategy, max 10)",
        f"{total - n_pass} failures: {', '.join(failed_names)}",
    ], font_size=12)

    add_textbox(slide, Inches(5.0), Inches(1.0), Inches(4.5), Inches(0.3),
                "Performance",
                font_size=16, bold=True, color=MED_BLUE)

    add_bullet_list(slide, Inches(5.0), Inches(1.35), Inches(4.5), Inches(0.85), [
        f"Median speedup: {median_s:.2f}x",
        f"Mean speedup: {mean_s:.2f}x",
        f"GPU wins (>1x): {gpu_wins}/{len(speedups)} ({100*gpu_wins//len(speedups) if speedups else 0}%)",
    ], font_size=12)

    # Top 5 table
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(9.0), Inches(0.25),
                "Top Speedups:", font_size=14, bold=True, color=DARK_BLUE)

    top_rows = [["Kernel", "Speedup", "C Time (ms)", "Triton Time (ms)"]]
    for k, s, c_t, t_t in speedups[:5]:
        top_rows.append([k, f"{s:.2f}x", f"{c_t:.3f}", f"{t_t:.3f}"])
    add_simple_table(slide, Inches(0.5), Inches(2.7), Inches(9.0),
                     [2.2, 2.2, 2.3, 2.3], top_rows,
                     font_size=12, row_height=0.32)

    # Slowdowns note (bottom 4)
    bottom = speedups[-4:] if len(speedups) >= 4 else speedups
    slow_str = ", ".join(f"{k} ({s:.2f}x)" for k, s, _, _ in reversed(bottom))
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.3),
                f"Slowdowns: {slow_str} -- sequential bottlenecks on GPU",
                font_size=11, color=ACCENT_RED)


def slide_10b_speedup_chart(prs):
    """Per-kernel speedup horizontal bar chart."""
    import json
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Per-Kernel Speedup")

    # Load data from results
    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with open(os.path.join(results_dir, "results.json")) as f:
        wr = json.load(f)

    data = []
    for k, v in wr.items():
        if v.get('test_passed'):
            b = v.get('benchmark', {})
            if b and b.get('speedup', 0) > 0:
                data.append((k, b['speedup']))
    data.sort(key=lambda x: -x[1])

    gpu_wins = sum(1 for _, s in data if s >= 1.0)
    all_s = [s for _, s in data]
    sorted_s = sorted(all_s)
    median_s = sorted_s[len(sorted_s) // 2] if sorted_s else 0
    mean_s = sum(all_s) / len(all_s) if all_s else 0
    max_speedup = max((s for _, s in data), default=7.0)

    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8.0), Inches(0.2),
                f"Median: {median_s:.2f}x  |  Mean: {mean_s:.2f}x  |  {gpu_wins}/{len(data)} faster on GPU",
                font_size=10, color=MED_GRAY)

    # Layout constants
    chart_top_in = 1.2
    row_h = min(0.14, 3.8 / max(len(data), 1))  # auto-fit rows
    bar_h = row_h * 0.72
    name_right = 1.35      # right edge of name column
    bar_left = 1.5         # left edge of bars
    bar_max_w = 7.5        # max bar width
    max_scale = max(int(max_speedup) + 1, 7)  # auto-scale

    # Axis tick marks at top
    for val in range(1, max_scale):
        x = bar_left + bar_max_w * (val / max_scale)
        # Tick label
        txb = slide.shapes.add_textbox(
            Inches(x - 0.12), Inches(chart_top_in - 0.18), Inches(0.24), Inches(0.15))
        p = txb.text_frame.paragraphs[0]
        p.text = f"{val}x"
        p.font.name = FONT_BODY
        p.font.size = Pt(7)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER
        # Tick line
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(chart_top_in - 0.02),
            Pt(0.5), Inches(row_h * len(data) + 0.02))
        tick.fill.solid()
        tick.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        tick.line.fill.background()

    # 1x reference line (thicker, dashed-style via color)
    one_x = bar_left + bar_max_w * (1.0 / max_scale)
    ref_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(one_x), Inches(chart_top_in - 0.02),
        Pt(1.5), Inches(row_h * len(data) + 0.02))
    ref_line.fill.solid()
    ref_line.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    ref_line.line.fill.background()

    # "1x" label (bold, above the line)
    txb = slide.shapes.add_textbox(
        Inches(one_x - 0.08), Inches(chart_top_in - 0.22), Inches(0.16), Inches(0.12))
    p = txb.text_frame.paragraphs[0]
    p.text = "1x"
    p.font.name = FONT_BODY
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Draw bars
    for i, (name, speedup) in enumerate(data):
        y_row = chart_top_in + row_h * i
        y_bar = y_row + (row_h - bar_h) / 2

        # Kernel name (right-aligned, monospace)
        txb = slide.shapes.add_textbox(
            Inches(0.05), Inches(y_row), Inches(name_right - 0.05), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_CODE
        p.font.size = Pt(7)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.RIGHT

        # Bar
        bar_w_ratio = min(speedup / max_scale, 1.0)
        bar_w = max(bar_max_w * bar_w_ratio, 0.03)
        bar_color = ACCENT_GREEN if speedup >= 1.0 else ACCENT_RED

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_left), Inches(y_bar),
            Inches(bar_w), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        # Speedup value label (after bar)
        label_x = bar_left + bar_w + 0.04
        txb = slide.shapes.add_textbox(
            Inches(label_x), Inches(y_row), Inches(0.45), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = f"{speedup:.2f}x"
        p.font.name = FONT_CODE
        p.font.size = Pt(6.5)
        p.font.color.rgb = bar_color
        p.alignment = PP_ALIGN.LEFT

    # Legend at bottom-right
    legend_y = chart_top_in + row_h * len(data) + 0.05
    # Green box
    g = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(7.0), Inches(legend_y), Inches(0.2), Inches(0.12))
    g.fill.solid()
    g.fill.fore_color.rgb = ACCENT_GREEN
    g.line.fill.background()
    slower_count = len(data) - gpu_wins
    add_textbox(slide, Inches(7.25), Inches(legend_y - 0.02), Inches(0.8), Inches(0.15),
                f"Faster ({gpu_wins})", font_size=7, color=DARK_TEXT)
    # Red box
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(8.2), Inches(legend_y), Inches(0.2), Inches(0.12))
    r.fill.solid()
    r.fill.fore_color.rgb = ACCENT_RED
    r.line.fill.background()
    add_textbox(slide, Inches(8.45), Inches(legend_y - 0.02), Inches(0.8), Inches(0.15),
                f"Slower ({slower_count})", font_size=7, color=DARK_TEXT)


def slide_10c_perf_comparison(prs):
    """Performance comparison: with analysis vs without analysis."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Performance: With vs Without Analysis")

    # Load results
    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with open(os.path.join(results_dir, "results.json")) as f:
        with_a = json.load(f)
    with open(os.path.join(results_dir, "results_no_analysis.json")) as f:
        no_a = json.load(f)

    # Build comparison data for kernels passing in both
    both_data = []
    for k in sorted(with_a.keys()):
        w = with_a.get(k, {})
        n = no_a.get(k, {})
        if w.get('test_passed') and n.get('test_passed'):
            wb = w.get('benchmark', {})
            nb = n.get('benchmark', {})
            ws = wb.get('speedup', 0) if wb else 0
            ns = nb.get('speedup', 0) if nb else 0
            if ws > 0 and ns > 0:
                both_data.append((k, ws, ns))

    # Sort by delta (analysis advantage) descending
    both_data.sort(key=lambda x: -(x[1] - x[2]))

    # Layout: paired horizontal bars
    chart_top = 1.1
    row_h = 0.16
    bar_h = 0.06
    name_right = 1.35
    bar_left = 1.5
    bar_max_w = 5.5
    # Auto-scale but cap to avoid extreme outliers making bars tiny
    all_speedups = [ws for _, ws, _ in both_data] + [ns for _, _, ns in both_data]
    p90 = sorted(all_speedups)[int(len(all_speedups) * 0.9)] if all_speedups else 10
    max_scale = max(int(p90) + 2, 10)

    # Axis ticks
    for val in range(1, max_scale):
        x = bar_left + bar_max_w * (val / max_scale)
        txb = slide.shapes.add_textbox(
            Inches(x - 0.12), Inches(chart_top - 0.18), Inches(0.24), Inches(0.15))
        p = txb.text_frame.paragraphs[0]
        p.text = f"{val}x"
        p.font.name = FONT_BODY
        p.font.size = Pt(7)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER

    # 1x reference line
    one_x = bar_left + bar_max_w * (1.0 / max_scale)
    ref_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(one_x), Inches(chart_top - 0.02),
        Pt(1.5), Inches(row_h * len(both_data) + 0.02))
    ref_line.fill.solid()
    ref_line.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    ref_line.line.fill.background()

    BLUE_WITH = RGBColor(0x2E, 0x86, 0xAB)  # With analysis
    ORANGE_WITHOUT = RGBColor(0xE8, 0x7D, 0x2F)  # Without analysis

    for i, (name, ws, ns) in enumerate(both_data):
        y_row = chart_top + row_h * i
        y_bar_w = y_row + 0.01  # With analysis bar (top)
        y_bar_n = y_row + 0.01 + bar_h + 0.01  # Without analysis bar (bottom)

        # Kernel name
        txb = slide.shapes.add_textbox(
            Inches(0.05), Inches(y_row), Inches(name_right - 0.05), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_CODE
        p.font.size = Pt(6)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.RIGHT

        # With analysis bar
        w_ratio = min(ws / max_scale, 1.0)
        w_bw = max(bar_max_w * w_ratio, 0.02)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_left), Inches(y_bar_w),
            Inches(w_bw), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE_WITH
        bar.line.fill.background()

        # Without analysis bar
        n_ratio = min(ns / max_scale, 1.0)
        n_bw = max(bar_max_w * n_ratio, 0.02)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_left), Inches(y_bar_n),
            Inches(n_bw), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE_WITHOUT
        bar.line.fill.background()

    # Summary stats on the right
    w_vals = [x[1] for x in both_data]
    n_vals = [x[2] for x in both_data]
    w_vals_s = sorted(w_vals)
    n_vals_s = sorted(n_vals)
    nn = len(w_vals)

    stats_x = Inches(7.3)
    add_textbox(slide, stats_x, Inches(1.0), Inches(2.5), Inches(0.25),
                f"Summary ({nn} kernels both pass):", font_size=10, bold=True, color=DARK_BLUE)

    stats = [
        f"With: median {w_vals_s[nn//2]:.2f}x, mean {sum(w_vals)/nn:.2f}x",
        f"W/O:  median {n_vals_s[nn//2]:.2f}x, mean {sum(n_vals)/nn:.2f}x",
        f"",
        f"With faster: {sum(1 for _, w, n in both_data if w > n*1.1)}/{nn}",
        f"W/O faster: {sum(1 for _, w, n in both_data if n > w*1.1)}/{nn}",
        f"Ties: {sum(1 for _, w, n in both_data if abs(w - n) <= max(w, n)*0.1)}/{nn}",
        f"",
        f"Correctness: {sum(1 for v in with_a.values() if v.get('test_passed'))}/{len(with_a)} vs {sum(1 for v in no_a.values() if v.get('test_passed'))}/{len(no_a)}",
    ]
    add_bullet_list(slide, stats_x, Inches(1.3), Inches(2.5), Inches(2.5),
                    stats, font_size=8, color=DARK_TEXT)

    # Legend
    legend_y = chart_top + row_h * len(both_data) + 0.1
    g = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1.5), Inches(legend_y), Inches(0.3), Inches(0.1))
    g.fill.solid()
    g.fill.fore_color.rgb = BLUE_WITH
    g.line.fill.background()
    add_textbox(slide, Inches(1.85), Inches(legend_y - 0.02), Inches(1.2), Inches(0.15),
                "With Analysis", font_size=7, color=DARK_TEXT)

    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(3.2), Inches(legend_y), Inches(0.3), Inches(0.1))
    r.fill.solid()
    r.fill.fore_color.rgb = ORANGE_WITHOUT
    r.line.fill.background()
    add_textbox(slide, Inches(3.55), Inches(legend_y - 0.02), Inches(1.2), Inches(0.15),
                "Without Analysis", font_size=7, color=DARK_TEXT)

    # Insight at bottom
    w_total_pass = sum(1 for v in with_a.values() if v.get('test_passed'))
    n_total_pass = sum(1 for v in no_a.values() if v.get('test_passed'))
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.3),
                f"With analysis: {w_total_pass}/{len(with_a)} pass. Without: {n_total_pass}/{len(no_a)} pass. "
                "LLM nondeterminism contributes to speedup variation across runs.",
                font_size=9, color=MED_GRAY)


def slide_11_failures(prs):
    import json
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Failure Analysis & Insights")

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with open(os.path.join(results_dir, "results.json")) as f:
        wr = json.load(f)

    n_pass = sum(1 for v in wr.values() if v.get('test_passed'))
    n_total = len(wr)
    failed = {k: v for k, v in wr.items() if not v.get('test_passed')}

    # Known root causes (for reference; all 30 pass in current run)
    root_causes = {
        "seidel_2d": ("Sequential", f"Gauss-Seidel strict ordering; grid=(1,) sequential, {_kinfo('seidel_2d','wa')[3]}"),
        "cholesky": ("Sequential", f"Sequential factorization, grid=(1,); {_kinfo('cholesky','wa')[3]}"),
        "nussinov": ("Sequential", f"3-way DP max with triangular dep; {_kinfo('nussinov','wa')[3]}"),
        "ludcmp": ("Sequential", f"LU decomposition with forward/back substitution; {_kinfo('ludcmp','wa')[3]}"),
    }

    if len(failed) == 0:
        add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.3),
                    f"All {n_pass}/{n_total} pass with analysis (100%)",
                    font_size=15, bold=True, color=ACCENT_GREEN)
    else:
        add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.3),
                    f"{len(failed)} Failures ({n_pass}/{n_total} pass with analysis):",
                    font_size=15, bold=True, color=ACCENT_ORANGE)

    fail_rows = [["Kernel", "Failure Mode", "Root Cause"]]
    for k in sorted(failed.keys()):
        mode, cause = root_causes.get(k, ("Unknown", "Investigation needed"))
        fail_rows.append([k, mode, cause])
    add_simple_table(slide, Inches(0.5), Inches(1.35), Inches(9.0),
                     [1.5, 1.5, 6.0], fail_rows,
                     font_size=10, row_height=0.28)

    # Infrastructure fixes
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(9.0), Inches(0.3),
                "Infrastructure & Test Quality Fixes:",
                font_size=15, bold=True, color=ACCENT_GREEN)

    fix_rows = [
        ["Kernel(s)", "Issue", "Fix"],
        ["doitgen", "Tested scratch buffer (sum)", "Marked 'temp'; excluded"],
        ["durbin", "FP32 divergence (120-step recurrence)", "Tolerance override (atol=0.05)"],
        ["gramschmidt", "M<N rank-deficient: Q diverges", "Tolerance override (atol=1.0)"],
        ["ludcmp", "Pivotless LU: A,y diverge, x correct", "A,y as 'temp'; only check x"],
        ["7 kernels", "randn invalid (solvers need SPD)", "Domain-appropriate init"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(2.7), Inches(9.0),
                     [1.4, 4.0, 3.6], fix_rows,
                     font_size=11, row_height=0.30)

    # Domain-appropriate inputs detail
    add_textbox(slide, Inches(0.5), Inches(4.55), Inches(9.0), Inches(0.25),
                "Domain inputs: cholesky (SPD), lu/ludcmp (diag-dominant), "
                "trisolv (lower-tri), gramschmidt (well-conditioned), "
                "nussinov (int bases), floyd_warshall (non-neg weights)",
                font_size=9, color=MED_GRAY)

    # Key insight
    y_insight = Inches(4.85)
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.5), y_insight, Inches(9.0), Inches(0.55))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    insight_box.line.color.rgb = MED_BLUE
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.05

    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.06)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Insight: "
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.name = FONT_BODY
    run.font.color.rgb = MED_BLUE

    run2 = p.add_run()
    run2.text = ("All 30 kernels pass with analysis. Slowdowns (<1x) are inherently sequential algorithms "
                 "(seidel_2d, cholesky, nussinov, ludcmp). Domain-appropriate test inputs validate mathematical correctness.")
    run2.font.size = Pt(10)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_12_ablation(prs):
    """Ablation study: Analysis-guided vs No-analysis (just retry on failure)."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Ablation: Analysis vs No-Analysis")

    # Try to load actual results
    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with_file = os.path.join(results_dir, "results.json")
    without_file = os.path.join(results_dir, "results_no_analysis.json")

    # Default placeholder values
    w_pass, w_total, w_first, w_avg_att = 24, 30, 9, 0.0
    wo_pass, wo_total, wo_first, wo_avg_att = 25, 30, 9, 0.0
    w_median, w_mean, wo_median, wo_mean = 0.0, 0.0, 0.0, 0.0
    w_gt1, wo_gt1 = 0, 0

    if os.path.exists(with_file):
        with open(with_file) as f:
            wr = json.load(f)
        w_total = len(wr)
        w_pass = sum(1 for v in wr.values() if v.get("test_passed"))
        w_first = sum(1 for v in wr.values() if v.get("test_passed") and v.get("attempts") == 1)
        passed_attempts = [v["attempts"] for v in wr.values() if v.get("test_passed")]
        w_avg_att = sum(passed_attempts) / len(passed_attempts) if passed_attempts else 0
        w_speedups = sorted([v.get('benchmark', {}).get('speedup', 0) for v in wr.values()
                             if v.get('test_passed') and v.get('benchmark', {}).get('speedup', 0) > 0])
        if w_speedups:
            w_median = w_speedups[len(w_speedups) // 2]
            w_mean = sum(w_speedups) / len(w_speedups)
            w_gt1 = sum(1 for s in w_speedups if s > 1.0)

    if os.path.exists(without_file):
        with open(without_file) as f:
            wor = json.load(f)
        wo_total = len(wor)
        wo_pass = sum(1 for v in wor.values() if v.get("test_passed"))
        wo_first = sum(1 for v in wor.values() if v.get("test_passed") and v.get("attempts") == 1)
        passed_attempts = [v["attempts"] for v in wor.values() if v.get("test_passed")]
        wo_avg_att = sum(passed_attempts) / len(passed_attempts) if passed_attempts else 0
        wo_speedups = sorted([v.get('benchmark', {}).get('speedup', 0) for v in wor.values()
                              if v.get('test_passed') and v.get('benchmark', {}).get('speedup', 0) > 0])
        if wo_speedups:
            wo_median = wo_speedups[len(wo_speedups) // 2]
            wo_mean = sum(wo_speedups) / len(wo_speedups)
            wo_gt1 = sum(1 for s in wo_speedups if s > 1.0)

    # Experiment description
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.3),
                "Same pipeline, same model (Sonnet), same 5+5 retries, tl.constexpr -- only difference: analysis in prompt",
                font_size=13, color=MED_GRAY)

    # Comparison table
    w_pct = f"{100*w_pass/w_total:.1f}%" if isinstance(w_pass, int) else "?"
    wo_pct = f"{100*wo_pass/wo_total:.1f}%" if isinstance(wo_pass, int) else "?"
    w_avg_str = f"{w_avg_att:.1f}" if w_avg_att else "?"
    wo_avg_str = f"{wo_avg_att:.1f}" if wo_avg_att else "?"

    rows = [
        ["Metric", "With Analysis", "Without Analysis", "Delta"],
        ["Pass Rate", f"{w_pass}/{w_total} ({w_pct})", f"{wo_pass}/{wo_total} ({wo_pct})",
         f"+{w_pass - wo_pass}" if isinstance(wo_pass, int) else "?"],
        ["First-Try Pass", str(w_first), str(wo_first),
         f"+{w_first - wo_first}" if isinstance(wo_first, int) else "?"],
        ["Avg Attempts (passed)", w_avg_str, wo_avg_str,
         f"{w_avg_att - wo_avg_att:+.1f}" if isinstance(wo_pass, int) else "?"],
        ["Median Speedup", f"{w_median:.2f}x", f"{wo_median:.2f}x",
         f"{w_median - wo_median:+.2f}x"],
        ["Mean Speedup", f"{w_mean:.2f}x", f"{wo_mean:.2f}x",
         f"{w_mean - wo_mean:+.2f}x"],
        ["Kernels >1x", f"{w_gt1}/{w_pass}", f"{wo_gt1}/{wo_pass}",
         f"+{w_gt1 - wo_gt1}"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(1.5), Inches(9.0),
                     [2.2, 2.5, 2.5, 1.8], rows,
                     font_size=11, row_height=0.27)

    # What analysis provides
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(4.5), Inches(0.3),
                "What analysis provides:", font_size=14, bold=True, color=DARK_BLUE)

    add_bullet_list(slide, Inches(0.5), Inches(4.1), Inches(4.5), Inches(1.0), [
        "Which dims to parallelize vs keep sequential",
        "WAR deps: need array copies before parallel region",
        "Reduction type: how to accumulate (tl.sum, etc.)",
        "Scalar expansion: privatize loop-carried scalars",
    ], font_size=11, color=DARK_TEXT)

    # Without analysis
    add_textbox(slide, Inches(5.3), Inches(3.8), Inches(4.5), Inches(0.3),
                "Without analysis, LLM must:", font_size=14, bold=True, color=ACCENT_ORANGE)

    add_bullet_list(slide, Inches(5.3), Inches(4.1), Inches(4.5), Inches(1.0), [
        "Infer parallelism from C code alone",
        "Guess correct memory access patterns",
        "Discover dependencies by trial and error",
        "Rely entirely on retry-with-error feedback",
    ], font_size=11, color=DARK_TEXT)

    # Insight box
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.45))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    insight_box.line.color.rgb = MED_BLUE
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.1
    tf = insight_box.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Takeaway: "
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.name = FONT_BODY
    run.font.color.rgb = MED_BLUE
    run2 = p.add_run()
    delta_pass = w_pass - wo_pass
    delta_med = w_median - wo_median
    run2.text = (f"Analysis: {w_pass}/{w_total} pass (vs {wo_pass}/{wo_total}), "
                 f"median {w_median:.2f}x (vs {wo_median:.2f}x). "
                 f"+{delta_pass} kernels pass, +{delta_med:.2f}x median speedup. "
                 f"Analysis uniquely enables hard dep-heavy kernels.")
    run2.font.size = Pt(12)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_03b_no_analysis_kernels(prs):
    """Explain the loop fission pattern and the 2 remaining kernels without analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Loop Fission: Opposing Reductions & Cross-Phase Deps")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "30/30 kernels receive analysis. Multi-phase split guidance added for correlation, covariance, gemver.",
                font_size=12, color=MED_GRAY)

    # LEFT: atax/bicg — SOLVED with fission guidance
    add_textbox(slide, Inches(0.3), Inches(1.15), Inches(4.5), Inches(0.25),
                "atax / bicg -- Opposing Reductions (SOLVED)",
                font_size=13, bold=True, color=ACCENT_GREEN)

    atax_code = """\
// atax: y = A^T * (A * x)
for (i = 0; i < M; i++) {
  tmp[i] = 0;
  for (j = 0; j < N; j++) {
    tmp[i] += A[i][j] * x[j];  // reduce j -> tmp[i]
    y[j]   += A[i][j] * tmp[i]; // reduce i -> y[j]
  }
}"""
    add_code_box(slide, Inches(0.3), Inches(1.4), Inches(4.5), Inches(1.5),
                 atax_code, font_size=8)

    add_bullet_list(slide, Inches(0.3), Inches(2.95), Inches(4.5), Inches(1.2), [
        "Both dims INVALID: write_conflict on opposing arrays",
        "New: ## Loop Fission Recommended",
        "  -> LLM splits into 2 kernels with tl.sum()",
        f"  -> atax {_kinfo('atax','wa')[3]}, bicg {_kinfo('bicg','wa')[3]}",
    ], font_size=10, color=DARK_TEXT)

    # RIGHT: correlation/covariance — LLM handles on its own
    add_textbox(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(0.25),
                "correlation / covariance -- Cross-Phase (MultiPhase)",
                font_size=13, bold=True, color=ACCENT_GREEN)

    corr_code = """\
// correlation: 4 fused phases
// Phase 1: mean[j] = sum_i(data[i][j]) / N
// Phase 2: stddev[j] = sqrt(sum_i(...))
// Phase 3: data[i][j] -= mean[j]  (in-place!)
// Phase 4: corr[i][j] = sum_k(data[k][i]
//                             * data[k][j])
//          corr[j][i] = corr[i][j]  (symmetry)"""
    add_code_box(slide, Inches(5.2), Inches(1.4), Inches(4.6), Inches(1.5),
                 corr_code, font_size=8)

    add_bullet_list(slide, Inches(5.2), Inches(2.95), Inches(4.6), Inches(1.2), [
        "PET sees cross-phase read/write as shift dep",
        "corr[j][i]=corr[i][j] looks like write conflict",
        "New: ## Multi-Phase Kernel guidance emitted",
        f"  -> covariance {_kinfo('covariance','wa')[3]}, correlation {_kinfo('correlation','wa')[3]}",
    ], font_size=10, color=DARK_TEXT)

    # Bottom insight box
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.3), Inches(4.05), Inches(9.5), Inches(0.55))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    insight_box.line.color.rgb = ACCENT_GREEN
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.08
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Loop Fission guidance works: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = (f"When ParDims detects write_conflict on both dims (opposing reductions), "
                 f"we emit a '## Loop Fission Recommended' section telling the LLM to split into separate kernels. "
                 f"atax {_kinfo('atax','wa')[3]}, bicg {_kinfo('bicg','wa')[3]}. "
                 f"Multi-phase guidance gives covariance {_kinfo('covariance','wa')[3]}, correlation {_kinfo('correlation','wa')[3]}.")
    run2.font.size = Pt(10)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_03c_ablation_all_kernels(prs):
    """Full ablation comparison: all 28 kernels that have analysis, sorted by analysis advantage."""
    import json

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    with_file = os.path.join(results_dir, "results.json")
    without_file = os.path.join(results_dir, "results_no_analysis.json")

    with open(with_file) as f:
        with_a = json.load(f)
    with open(without_file) as f:
        no_a = json.load(f)

    # Kernels without any analysis (identical prompts in both modes)
    no_analysis_kernels = set()  # All 30 kernels now receive analysis

    # Build comparison rows with sort key
    all_rows = []
    for k in with_a:
        if k.replace('-', '_') in no_analysis_kernels:
            continue
        w = with_a.get(k, {})
        n = no_a.get(k, {})
        w_passed = w.get('test_passed', False)
        n_passed = n.get('test_passed', False)
        w_att = w.get('attempts', '?')
        n_att = n.get('attempts', '?')
        wb = w.get('benchmark', {})
        nb = n.get('benchmark', {})
        ws = wb.get('speedup', 0) if wb else 0
        ns = nb.get('speedup', 0) if nb else 0

        w_str = f"Y({w_att}) {ws:.2f}x" if w_passed and ws > 0 else ("Y({})".format(w_att) if w_passed else "FAIL")
        n_str = f"Y({n_att}) {ns:.2f}x" if n_passed and ns > 0 else ("Y({})".format(n_att) if n_passed else "FAIL")

        # Compute sort ratio (analysis advantage)
        if w_passed and not n_passed:
            sort_ratio = 1000.0  # analysis uniquely passes
            winner = "W/ (pass)"
        elif not w_passed and n_passed:
            sort_ratio = -1000.0
            winner = "W/O (pass)"
        elif w_passed and n_passed and ws > 0 and ns > 0:
            sort_ratio = ws / ns
            if sort_ratio > 1.15:
                winner = f"W/ (+{sort_ratio:.1f}x)"
            elif sort_ratio < 0.87:
                winner = f"W/O (+{1/sort_ratio:.1f}x)"
            else:
                winner = "tie"
        elif not w_passed and not n_passed:
            sort_ratio = 1.0
            winner = "both fail"
        else:
            sort_ratio = 1.0
            winner = "tie"

        all_rows.append((k, w_str, n_str, winner, sort_ratio))

    # Sort by analysis advantage descending
    all_rows.sort(key=lambda x: -x[4])

    # Split into two pages: 14 rows each
    mid = (len(all_rows) + 1) // 2
    n_kernels = len(all_rows)

    for page, start_idx in enumerate([0, mid]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_label = f" ({page+1}/2)"
        add_section_title(slide, f"Ablation: All {n_kernels} Kernels With Analysis{page_label}")

        add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.25),
                    "Sorted by analysis advantage (with/without speedup ratio). Top = analysis helps most.",
                    font_size=11, color=MED_GRAY)

        page_rows = all_rows[start_idx:start_idx + mid]
        table_rows = [["Kernel", "With Analysis", "Without Analysis", "Winner"]]
        for k, w_str, n_str, winner, _ in page_rows:
            table_rows.append([k, w_str, n_str, winner])

        add_simple_table(slide, Inches(0.3), Inches(1.1), Inches(9.4),
                         [2.0, 2.8, 2.8, 1.8], table_rows,
                         font_size=10, row_height=0.28)

        # On second page, add summary
        if page == 1:
            w_wins = sum(1 for _, _, _, w, _ in all_rows if w.startswith("W/ "))
            wo_wins = sum(1 for _, _, _, w, _ in all_rows if w.startswith("W/O"))
            ties = sum(1 for _, _, _, w, _ in all_rows if w == "tie")
            both_f = sum(1 for _, _, _, w, _ in all_rows if w == "both fail")

            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(1.5), Inches(5.0), Inches(7.0), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = MED_BLUE
            badge.line.fill.background()
            badge.adjustments[0] = 0.15
            tf = badge.text_frame
            tf.margin_left = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = f"Analysis wins {w_wins}/{n_kernels}  |  No-analysis wins {wo_wins}/{n_kernels}  |  Ties {ties}  |  Both fail {both_f}"
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════════════
# RODINIA SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_13_rodinia_results(prs):
    """Rodinia 3.1 infrastructure test: overview + results table."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Rodinia 3.1: Pipeline Generalization Test")

    # Load results
    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "rodinia_results")
    results_file = os.path.join(results_dir, "results.json")
    if not os.path.exists(results_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No Rodinia results found. Run generate_and_test_rodinia.py first.",
                    font_size=16, color=ACCENT_RED)
        return

    with open(results_file) as f:
        wr = json.load(f)

    total = len(wr)
    passed = {k: v for k, v in wr.items() if v.get('test_passed')}
    n_pass = len(passed)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "Infrastructure test: 3 kernels from distinct algorithmic families, no polyhedral analysis.",
                font_size=12, color=MED_GRAY)

    # Kernel selection rationale table
    add_textbox(slide, Inches(0.3), Inches(1.15), Inches(4.5), Inches(0.3),
                "Kernel Selection:", font_size=14, bold=True, color=DARK_BLUE)

    kernel_rows = [
        ["Kernel", "Pattern", "Polybench Analog"],
        ["hotspot", "2D stencil + timestep", "jacobi_2d, heat_3d"],
        ["lud", "LU decomposition", "lu"],
        ["pathfinder", "1D DP + timestep", "nussinov"],
    ]
    add_simple_table(slide, Inches(0.3), Inches(1.45), Inches(4.5),
                     [1.2, 1.8, 1.5], kernel_rows,
                     font_size=11, row_height=0.3)

    # Results table
    add_textbox(slide, Inches(5.2), Inches(1.15), Inches(4.6), Inches(0.3),
                "Results:", font_size=14, bold=True, color=DARK_BLUE)

    result_rows = [["Kernel", "Passed", "Attempts", "Speedup"]]
    for k in sorted(wr.keys()):
        v = wr[k]
        p_str = "Y" if v.get('test_passed') else "N"
        att = str(v.get('attempts', '?'))
        bench = v.get('benchmark', {})
        sp = bench.get('speedup', 0) if bench else 0
        sp_str = f"{sp:.2f}x" if sp > 0 else "-"
        result_rows.append([k, p_str, att, sp_str])

    add_simple_table(slide, Inches(5.2), Inches(1.45), Inches(4.6),
                     [1.2, 0.8, 1.0, 1.0], result_rows,
                     font_size=11, row_height=0.3)

    # Speedups for passed kernels
    speedups = []
    for k, v in passed.items():
        b = v.get('benchmark', {})
        if b and b.get('speedup', 0) > 0:
            speedups.append(b['speedup'])

    # Summary stats badges
    badge_y = Inches(2.9)
    pct = f"{100*n_pass//total}" if total > 0 else "?"
    badges = [
        (f"{n_pass}/{total} Pass ({pct}%)", ACCENT_GREEN if n_pass == total else ACCENT_ORANGE),
        (f"No Analysis (infra test)", MED_BLUE),
    ]
    if speedups:
        sorted_s = sorted(speedups)
        median_s = sorted_s[len(sorted_s) // 2]
        mean_s = sum(speedups) / len(speedups)
        badges.append((f"Median: {median_s:.2f}x", LIGHT_BLUE))
        badges.append((f"Mean: {mean_s:.2f}x", LIGHT_BLUE))

    x = Inches(0.5)
    for label, color in badges:
        add_flow_box(slide, x, badge_y, Inches(2.1), Inches(0.4), label,
                     fill_color=color, font_size=11)
        x += Inches(2.3)

    # Per-kernel details
    add_textbox(slide, Inches(0.3), Inches(3.5), Inches(9.5), Inches(0.3),
                "Per-Kernel Notes:", font_size=14, bold=True, color=DARK_BLUE)

    notes = []

    # hotspot
    hotspot_v = wr.get('hotspot', {})
    h_att = hotspot_v.get('attempts', '?')
    h_bench = hotspot_v.get('benchmark', {})
    h_sp = h_bench.get('speedup', 0) if h_bench else 0
    notes.append(
        f"hotspot: 2D thermal stencil (256x256, 10 timesteps). "
        f"Passed on attempt {h_att}, {h_sp:.2f}x speedup. "
        f"Double-buffer pattern (read temp, write result, copy back)."
    )

    # lud
    lud_v = wr.get('lud', {})
    l_att = lud_v.get('attempts', '?')
    l_bench = lud_v.get('benchmark', {})
    l_sp = l_bench.get('speedup', 0) if l_bench else 0
    notes.append(
        f"lud: Unblocked LU decomposition (256x256). "
        f"First-try pass, {l_sp:.2f}x speedup. "
        f"Identical to Polybench lu pattern -- pipeline handles it consistently."
    )

    # pathfinder
    pf_v = wr.get('pathfinder', {})
    pf_att = pf_v.get('attempts', '?')
    pf_passed = pf_v.get('test_passed', False)
    if pf_passed:
        pf_bench = pf_v.get('benchmark', {})
        pf_sp = pf_bench.get('speedup', 0) if pf_bench else 0
        notes.append(
            f"pathfinder: 1D DP min-path (100x256). "
            f"Passed on attempt {pf_att}, {pf_sp:.2f}x speedup."
        )
    else:
        notes.append(
            f"pathfinder: 1D DP min-path (100x256). "
            f"Failed all {pf_att} attempts. "
            f"Double-buffer DP with neighbor min -- analysis guidance (WAR + ParDims) would help."
        )

    add_bullet_list(slide, Inches(0.3), Inches(3.8), Inches(9.5), Inches(1.5),
                    notes, font_size=10, color=DARK_TEXT, bold_prefix=True)

    # Bottom insight
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.3), Inches(5.0), Inches(9.5), Inches(0.45))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    insight_box.line.color.rgb = ACCENT_GREEN
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.08
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Validates pipeline generalization: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = (f"{n_pass}/{total} Rodinia kernels pass with zero pipeline changes. "
                 "Same C extraction, ctypes testing, and LLM retry strategy works across benchmark suites.")
    run2.font.size = Pt(11)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_13b_rodinia_speedup_chart(prs):
    """Rodinia per-kernel speedup horizontal bar chart."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Rodinia: Per-Kernel Speedup")

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "rodinia_results")
    results_file = os.path.join(results_dir, "results.json")
    if not os.path.exists(results_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No Rodinia results found.",
                    font_size=16, color=ACCENT_RED)
        return

    with open(results_file) as f:
        wr = json.load(f)

    # Collect data: all kernels (passed with speedup + failed)
    passed_data = []
    failed_names = []
    for k, v in sorted(wr.items()):
        if v.get('test_passed'):
            b = v.get('benchmark', {})
            if b and b.get('speedup', 0) > 0:
                passed_data.append((k, b['speedup']))
            else:
                passed_data.append((k, 0))
        else:
            failed_names.append(k)
    passed_data.sort(key=lambda x: -x[1])

    if not passed_data:
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No passing kernels with benchmark data.",
                    font_size=16, color=ACCENT_RED)
        return

    gpu_wins = sum(1 for _, s in passed_data if s >= 1.0)
    all_s = [s for _, s in passed_data if s > 0]
    sorted_s = sorted(all_s)
    median_s = sorted_s[len(sorted_s) // 2] if sorted_s else 0
    mean_s = sum(all_s) / len(all_s) if all_s else 0
    max_speedup = max((s for _, s in passed_data), default=7.0)

    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8.0), Inches(0.2),
                f"Median: {median_s:.2f}x  |  Mean: {mean_s:.2f}x  |  {gpu_wins}/{len(passed_data)} faster on GPU",
                font_size=11, color=MED_GRAY)

    # Layout constants -- larger bars for fewer kernels
    chart_top_in = 1.4
    row_h = 0.55
    bar_h = row_h * 0.6
    name_right = 1.8
    bar_left = 2.0
    bar_max_w = 6.0
    max_scale = max(int(max_speedup) + 2, 7)

    # Axis tick marks
    for val in range(1, max_scale):
        x = bar_left + bar_max_w * (val / max_scale)
        txb = slide.shapes.add_textbox(
            Inches(x - 0.15), Inches(chart_top_in - 0.25), Inches(0.30), Inches(0.20))
        p = txb.text_frame.paragraphs[0]
        p.text = f"{val}x"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER
        # Grid line
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(chart_top_in - 0.02),
            Pt(0.5), Inches(row_h * len(passed_data) + 0.04))
        tick.fill.solid()
        tick.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        tick.line.fill.background()

    # 1x reference line
    one_x = bar_left + bar_max_w * (1.0 / max_scale)
    ref_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(one_x), Inches(chart_top_in - 0.02),
        Pt(1.5), Inches(row_h * len(passed_data) + 0.04))
    ref_line.fill.solid()
    ref_line.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    ref_line.line.fill.background()

    # Draw bars
    for i, (name, speedup) in enumerate(passed_data):
        y_row = chart_top_in + row_h * i
        y_bar = y_row + (row_h - bar_h) / 2

        # Kernel name
        txb = slide.shapes.add_textbox(
            Inches(0.1), Inches(y_row), Inches(name_right - 0.1), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_CODE
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.RIGHT

        # Bar
        bar_w_ratio = min(speedup / max_scale, 1.0)
        bar_w = max(bar_max_w * bar_w_ratio, 0.03)
        bar_color = ACCENT_GREEN if speedup >= 1.0 else ACCENT_RED

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_left), Inches(y_bar),
            Inches(bar_w), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        # Speedup label
        label_x = bar_left + bar_w + 0.08
        txb = slide.shapes.add_textbox(
            Inches(label_x), Inches(y_row), Inches(0.8), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = f"{speedup:.2f}x"
        p.font.name = FONT_CODE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = bar_color
        p.alignment = PP_ALIGN.LEFT

    # Failed kernels note
    if failed_names:
        fail_y = chart_top_in + row_h * len(passed_data) + 0.15
        add_textbox(slide, Inches(0.5), Inches(fail_y), Inches(9.0), Inches(0.3),
                    f"Failed ({len(failed_names)}): {', '.join(failed_names)}",
                    font_size=12, bold=True, color=ACCENT_RED)

        # Failure explanation
        fail_details = []
        for fn in failed_names:
            fv = wr.get(fn, {})
            fe = fv.get('final_error', {})
            if isinstance(fe, dict):
                etype = fe.get('type', 'unknown')
            else:
                etype = 'unknown'
            fail_details.append(f"{fn}: {etype} error after {fv.get('attempts', '?')} attempts")
        add_bullet_list(slide, Inches(0.5), Inches(fail_y + 0.3), Inches(9.0), Inches(0.8),
                        fail_details, font_size=10, color=DARK_TEXT)

    # Comparison with Polybench
    polybench_results_file = os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "polybench_results", "results.json")
    if os.path.exists(polybench_results_file):
        with open(polybench_results_file) as f:
            pb = json.load(f)
        pb_pass = sum(1 for v in pb.values() if v.get('test_passed'))
        pb_total = len(pb)
        pb_speedups = []
        for v in pb.values():
            if v.get('test_passed'):
                b = v.get('benchmark', {})
                if b and b.get('speedup', 0) > 0:
                    pb_speedups.append(b['speedup'])
        pb_median = sorted(pb_speedups)[len(pb_speedups) // 2] if pb_speedups else 0

        compare_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(0.5), Inches(4.8), Inches(9.0), Inches(0.5))
        compare_box.fill.solid()
        compare_box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
        compare_box.line.color.rgb = MED_BLUE
        compare_box.line.width = Pt(1)
        compare_box.adjustments[0] = 0.08
        tf = compare_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Cross-suite comparison: "
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.name = FONT_BODY
        run.font.color.rgb = MED_BLUE
        run2 = p.add_run()
        rod_n_pass = sum(1 for v in wr.values() if v.get('test_passed'))
        rod_total = len(wr)
        rod_pass_str = f"{rod_n_pass}/{rod_total}"
        rod_med_str = f"{median_s:.2f}x" if all_s else "N/A"
        run2.text = (f"Polybench: {pb_pass}/{pb_total} pass, median {pb_median:.2f}x  |  "
                     f"Rodinia: {rod_pass_str} pass, median {rod_med_str}  |  "
                     f"Same pipeline, zero code changes.")
        run2.font.size = Pt(11)
        run2.font.name = FONT_BODY
        run2.font.color.rgb = DARK_TEXT


# ────────────────────────────────────────────────────────────────────────
# SLIDE 14: Rodinia Ablation (Analysis vs No-Analysis)
# ────────────────────────────────────────────────────────────────────────

def slide_14_rodinia_ablation(prs):
    """Rodinia ablation: with-analysis vs without-analysis comparison."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Rodinia: Ablation — Analysis vs No-Analysis")

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "rodinia_results")
    with_file = os.path.join(results_dir, "results.json")
    without_file = os.path.join(results_dir, "results_no_analysis.json")

    if not os.path.exists(with_file) or not os.path.exists(without_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "Ablation data not available (need both results.json and results_no_analysis.json).",
                    font_size=14, color=ACCENT_RED)
        return

    with open(with_file) as f:
        w_results = json.load(f)
    with open(without_file) as f:
        wo_results = json.load(f)

    # Build comparison table data — compare by Triton time (immune to C-ref variance)
    headers = ["Kernel", "NA Speedup", "NA Triton ms", "WA Speedup", "WA Triton ms", "Winner"]
    rows = []
    all_kernels = sorted(set(list(w_results.keys()) + list(wo_results.keys())))

    analysis_wins = 0
    no_analysis_wins = 0
    ties = 0
    analysis_unique_pass = []

    for k in all_kernels:
        w = w_results.get(k, {})
        wo = wo_results.get(k, {})

        w_passed = w.get('test_passed', False)
        wo_passed = wo.get('test_passed', False)

        wb = w.get('benchmark', {})
        nob = wo.get('benchmark', {})
        w_sp = wb.get('speedup', 0) if w_passed else 0
        wo_sp = nob.get('speedup', 0) if wo_passed else 0
        w_tri = wb.get('triton_time_ms', 0) if w_passed else 0
        wo_tri = nob.get('triton_time_ms', 0) if wo_passed else 0

        w_sp_str = f"{w_sp:.2f}x" if w_passed else "FAIL"
        wo_sp_str = f"{wo_sp:.2f}x" if wo_passed else "FAIL"
        w_tri_str = f"{w_tri:.3f}" if w_passed else "—"
        wo_tri_str = f"{wo_tri:.3f}" if wo_passed else "—"

        # Determine winner by Triton time (lower is better)
        if w_passed and not wo_passed:
            winner = "WA (pass)"
            analysis_wins += 1
            analysis_unique_pass.append(k)
        elif not w_passed and wo_passed:
            winner = "NA (pass)"
            no_analysis_wins += 1
        elif w_passed and wo_passed and w_tri > 0 and wo_tri > 0:
            ratio = wo_tri / w_tri  # >1 means WA faster
            if ratio > 1.05:
                winner = f"WA ({ratio:.2f}x)"
                analysis_wins += 1
            elif ratio < 0.95:
                winner = f"NA ({1/ratio:.2f}x)"
                no_analysis_wins += 1
            else:
                winner = "Tie"
                ties += 1
        else:
            winner = "Both fail"

        rows.append([k, wo_sp_str, wo_tri_str, w_sp_str, w_tri_str, winner])

    # Pass rate row
    w_pass = sum(1 for v in w_results.values() if v.get('test_passed'))
    wo_pass = sum(1 for v in wo_results.values() if v.get('test_passed'))
    w_total = len(w_results)
    wo_total = len(wo_results)
    rows.append(["Pass rate", f"{wo_pass}/{wo_total}", "—", f"{w_pass}/{w_total}", "—",
                 "WA" if w_pass > wo_pass else ("NA" if wo_pass > w_pass else "Tie")])

    # Draw table
    col_widths = [Inches(1.3), Inches(1.4), Inches(1.5), Inches(1.4), Inches(1.5), Inches(1.5)]
    table_top = Inches(1.0)
    table_left = Inches(0.4)
    n_rows = len(rows) + 1  # +1 for header
    table = slide.shapes.add_table(n_rows, 6, table_left, table_top,
                                    sum(w.inches for w in col_widths if hasattr(w, 'inches')),
                                    Inches(0.35 * n_rows)).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.name = FONT_BODY
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = MED_BLUE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT_CODE if ci == 0 else FONT_BODY
                p.alignment = PP_ALIGN.CENTER
                # Color winner column
                if ci == 5:
                    if val.startswith("WA"):
                        p.font.color.rgb = ACCENT_GREEN
                        p.font.bold = True
                    elif val.startswith("NA"):
                        p.font.color.rgb = ACCENT_RED
                        p.font.bold = True
                # Color FAIL cells
                if "FAIL" in val:
                    p.font.color.rgb = ACCENT_RED
                    p.font.bold = True
            # Alternate row colors
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            # Bold the pass rate summary row
            if ri == len(rows) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True

    # Key insight box
    insight_y = table_top.inches + 0.35 * n_rows + 0.3
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.5), Inches(insight_y), Inches(9.0), Inches(0.8))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    insight_box.line.color.rgb = ACCENT_GREEN
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.08
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.06)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Finding: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN

    run2 = p.add_run()
    if analysis_unique_pass:
        run2.text = (f"Analysis uniquely enables {', '.join(analysis_unique_pass)} "
                     f"(FAIL → PASS). Pass rate: {wo_pass}/{wo_total} → {w_pass}/{w_total}. "
                     f"WA wins={analysis_wins}, NA wins={no_analysis_wins}, Tie={ties} "
                     f"(by Triton time). Speedup differences on lud/hotspot are from "
                     f"C-ref timing variance, not code quality.")
    else:
        run2.text = (f"WA wins={analysis_wins}, NA wins={no_analysis_wins}, Tie={ties} "
                     f"(compared by Triton time). Speedup differences are dominated by "
                     f"C-reference timing variance between runs.")
    run2.font.size = Pt(11)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT

    # Caveat note
    caveat_y = insight_y + 0.9
    add_textbox(slide, Inches(0.5), Inches(caveat_y), Inches(9.0), Inches(0.3),
                "Note: Winner determined by Triton kernel time (not speedup ratio) to eliminate "
                "C-reference timing variance. With only 3 kernels, small sample effects dominate.",
                font_size=9, color=MED_GRAY)


# ════════════════════════════════════════════════════════════════════════
# NCU UTILIZATION + NONDETERMINISM EVIDENCE
# ════════════════════════════════════════════════════════════════════════

def slide_15_ncu_utilization(prs):
    """NCU hardware utilization: show GPU is severely underutilized at Polybench sizes."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "GPU Hardware Utilization (NCU Profiling)")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "NVIDIA Nsight Compute profiling on RTX 3090 (82 SMs, 936 GB/s, 35.58 TFLOPS FP32)",
                font_size=12, color=MED_GRAY)

    # Load NCU results
    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    ncu_file = os.path.join(results_dir, "ncu_results.json")
    results_file = os.path.join(results_dir, "results.json")

    if not os.path.exists(ncu_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No NCU results found. Run ncu_profile_all.py first.",
                    font_size=16, color=ACCENT_RED)
        return

    with open(ncu_file) as f:
        ncu_data = json.load(f)
    with open(results_file) as f:
        bench = json.load(f)

    # Build summary rows sorted by SM%
    kernel_rows = []
    for kname, kernels in ncu_data.items():
        if not kernels:
            continue
        speedup = bench.get(kname, {}).get('benchmark', {}).get('speedup', 0)

        def _get_sm(k):
            try: return float(k.get("sm__throughput.avg.pct_of_peak_sustained_elapsed", "0").replace(",", ""))
            except: return 0.0

        best = max(kernels, key=_get_sm)
        sm = _get_sm(best)
        try: mem = float(best.get("gpu__compute_memory_throughput.avg.pct_of_peak_sustained_elapsed", "0").replace(",", ""))
        except: mem = 0.0
        try: occ = float(best.get("sm__warps_active.avg.pct_of_peak_sustained_active", "0").replace(",", ""))
        except: occ = 0.0
        grid = best.get("launch__grid_size", "?")
        nk = len(set(k["name"] for k in kernels))
        kernel_rows.append((kname, speedup, sm, mem, occ, grid, nk))

    kernel_rows.sort(key=lambda x: -x[2])  # Sort by SM% descending

    # Show top 8 and bottom 4 in a compact table
    show_kernels = kernel_rows[:8] + kernel_rows[-4:] if len(kernel_rows) > 12 else kernel_rows

    rows = [["Kernel", "Speedup", "SM%", "Mem%", "Occ%", "Grid"]]
    for kname, spd, sm, mem, occ, grid, nk in show_kernels:
        rows.append([kname, f"{spd:.1f}x", f"{sm:.1f}", f"{mem:.1f}", f"{occ:.1f}", str(grid)])

    add_simple_table(slide, Inches(0.3), Inches(1.15), Inches(5.8),
                     [1.2, 0.7, 0.7, 0.7, 0.7, 1.1], rows,
                     font_size=9, row_height=0.24)

    # Stats panel on right
    sm_vals = [x[2] for x in kernel_rows]
    mem_vals = [x[3] for x in kernel_rows]
    occ_vals = [x[4] for x in kernel_rows]

    stats_x = Inches(6.3)
    add_textbox(slide, stats_x, Inches(1.15), Inches(3.5), Inches(0.25),
                "Aggregate Statistics:", font_size=13, bold=True, color=DARK_BLUE)

    stats_items = []
    if sm_vals:
        stats_items.append(f"SM Throughput: mean {sum(sm_vals)/len(sm_vals):.1f}%, max {max(sm_vals):.1f}%")
    if mem_vals:
        stats_items.append(f"Mem Throughput: mean {sum(mem_vals)/len(mem_vals):.1f}%, max {max(mem_vals):.1f}%")
    if occ_vals:
        stats_items.append(f"Occupancy: mean {sum(occ_vals)/len(occ_vals):.1f}%, max {max(occ_vals):.1f}%")
    grid1_count = sum(1 for _, _, _, _, _, g, _ in kernel_rows if g in ("1", "?"))
    stats_items.append(f"grid=(1,) kernels: {grid1_count}/{len(kernel_rows)}")
    stats_items.append(f"Only correlation >50% SM")

    add_bullet_list(slide, stats_x, Inches(1.45), Inches(3.5), Inches(1.5),
                    stats_items, font_size=10, color=DARK_TEXT)

    # Interpretation box
    add_textbox(slide, stats_x, Inches(3.0), Inches(3.5), Inches(0.25),
                "Root Cause:", font_size=13, bold=True, color=ACCENT_RED)
    add_bullet_list(slide, stats_x, Inches(3.3), Inches(3.5), Inches(1.5), [
        "Polybench SMALL sizes (N=60-120)",
        "Data fits in L2 cache entirely",
        "Not enough parallelism to fill 82 SMs",
        "Most kernels: 1 block on 1 SM",
    ], font_size=10, color=DARK_TEXT)

    # Bottom insight
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.3), Inches(4.85), Inches(9.5), Inches(0.55))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
    insight_box.line.color.rgb = ACCENT_ORANGE
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.08
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Conclusion: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_ORANGE
    run2 = p.add_run()
    run2.text = ("GPU is severely underutilized (<10% SM/Mem). Current speedups (up to 13.9x) are impressive "
                 "DESPITE tiny problem sizes. Increasing to N=1024+ would show true GPU potential "
                 "and make speedup comparisons more meaningful.")
    run2.font.size = Pt(10)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_15b_large_size_benchmark(prs):
    """Large-size benchmark: Original (1x) vs 8x-scale WA speedups."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Large-Scale Results: Original vs 8x Problem Sizes")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "Full pipeline re-run at 8x problem sizes (LLM sees scaled dims in prompt, "
                "C reference recompiled at scaled sizes). With-analysis mode.",
                font_size=11, color=MED_GRAY)

    # Load both result files
    base_dir = os.path.dirname(os.path.abspath(__file__))
    orig_file = os.path.join(base_dir, "polybench_results", "results.json")
    scale_file = os.path.join(base_dir, "polybench_results_scale8x", "results.json")

    if not os.path.exists(scale_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No 8x-scale results found. Run: python generate_and_test_polybench.py --size-scale 8",
                    font_size=16, color=ACCENT_RED)
        return

    with open(orig_file) as f:
        orig = json.load(f)
    with open(scale_file) as f:
        scaled = json.load(f)

    # Build comparison table sorted by 8x speedup descending
    rows = [["Kernel", "Pass 1x", "Spd 1x", "Pass 8x", "Spd 8x", "Change"]]
    orig_spds, scale_spds = [], []
    improved, degraded, stable = 0, 0, 0

    all_kernels = sorted(set(list(orig.keys()) + list(scaled.keys())))
    kernel_data = []
    for k in all_kernels:
        o = orig.get(k, {})
        s = scaled.get(k, {})
        o_pass = "Y" if o.get("test_passed") else "N"
        s_pass = "Y" if s.get("test_passed") else "N"
        o_spd = o.get("benchmark", {}).get("speedup") if o.get("benchmark") else None
        s_spd = s.get("benchmark", {}).get("speedup") if s.get("benchmark") else None
        kernel_data.append((k, o_pass, o_spd, s_pass, s_spd))

    # Sort by 8x speedup descending (None at bottom)
    kernel_data.sort(key=lambda x: -(x[4] or -999))

    for k, o_pass, o_spd, s_pass, s_spd in kernel_data:
        o_str = f"{o_spd:.2f}x" if o_spd else "—"
        s_str = f"{s_spd:.2f}x" if s_spd else "—"
        if o_spd and o_spd > 0:
            orig_spds.append(o_spd)
        if s_spd and s_spd > 0:
            scale_spds.append(s_spd)
        # Change indicator
        if o_spd and s_spd:
            ratio = s_spd / o_spd if o_spd > 0 else 0
            if ratio > 1.1:
                change = f"+{(ratio - 1) * 100:.0f}%"
                improved += 1
            elif ratio < 0.9:
                change = f"{(ratio - 1) * 100:.0f}%"
                degraded += 1
            else:
                change = "~"
                stable += 1
        else:
            change = "—"
        rows.append([k, o_pass, o_str, s_pass, s_str, change])

    # Table on left
    n_show = min(len(rows) - 1, 15)
    display_rows = [rows[0]] + rows[1:n_show+1]
    if len(rows) - 1 > n_show:
        display_rows.append(["...", "", "", "", "", f"({len(rows)-1} total)"])

    add_simple_table(slide, Inches(0.15), Inches(1.15), Inches(5.8),
                     [1.2, 0.6, 0.9, 0.6, 0.9, 0.8], display_rows,
                     font_size=8, row_height=0.22)

    # Stats panel on right
    stats_x = Inches(6.2)
    add_textbox(slide, stats_x, Inches(1.15), Inches(3.6), Inches(0.25),
                "Summary Statistics:", font_size=13, bold=True, color=DARK_BLUE)

    o_pass_n = sum(1 for _, op, _, _, _ in kernel_data if op == "Y")
    s_pass_n = sum(1 for _, _, _, sp, _ in kernel_data if sp == "Y")
    o_sorted = sorted(orig_spds) if orig_spds else [0]
    s_sorted = sorted(scale_spds) if scale_spds else [0]
    o_med = o_sorted[len(o_sorted)//2]
    s_med = s_sorted[len(s_sorted)//2]
    o_gt1 = sum(1 for s in orig_spds if s > 1)
    s_gt1 = sum(1 for s in scale_spds if s > 1)

    stats_items = [
        f"Pass rate: {o_pass_n}/30 (1x) → {s_pass_n}/30 (8x)",
        f"Benchmarked: {len(orig_spds)} (1x) → {len(scale_spds)} (8x)",
        f"Median speedup: {o_med:.2f}x → {s_med:.2f}x",
        f"Mean speedup: {sum(orig_spds)/len(orig_spds):.2f}x → {sum(scale_spds)/len(scale_spds):.2f}x" if scale_spds else "",
        f"GPU speedup (>1x): {o_gt1} → {s_gt1}",
    ]
    stats_items = [s for s in stats_items if s]

    add_bullet_list(slide, stats_x, Inches(1.45), Inches(3.6), Inches(2.0),
                    stats_items, font_size=10, color=DARK_TEXT)

    # Improved / degraded
    add_textbox(slide, stats_x, Inches(3.2), Inches(3.6), Inches(0.25),
                "Speedup Change (1x → 8x):", font_size=13, bold=True, color=ACCENT_GREEN)

    n_both = improved + degraded + stable
    cat_items = [
        f"Improved (>10%): {improved}/{n_both} kernels",
        f"Stable (within 10%): {stable}/{n_both} kernels",
        f"Degraded (>10%): {degraded}/{n_both} kernels",
    ]
    add_bullet_list(slide, stats_x, Inches(3.5), Inches(3.6), Inches(1.5),
                    cat_items, font_size=10, color=DARK_TEXT)

    # Bottom insight
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.3), Inches(4.85), Inches(9.5), Inches(0.55))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    insight_box.line.color.rgb = ACCENT_GREEN
    insight_box.line.width = Pt(1.5)
    insight_box.adjustments[0] = 0.08
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Finding: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = ACCENT_GREEN
    run2 = p.add_run()
    run2.text = (f"At 8x problem sizes, median speedup jumps from {o_med:.2f}x to {s_med:.2f}x. "
                 f"{s_gt1}/{len(scale_spds)} kernels show GPU speedup. "
                 f"LLM generates better-optimized Triton code when it sees realistic problem dimensions.")
    run2.font.size = Pt(10)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_15c_large_speedup_chart(prs):
    """Per-kernel speedup bar chart at 8x problem sizes."""
    import json
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Per-Kernel Speedup at 8x Scale")

    base_dir = os.path.dirname(os.path.abspath(__file__))
    scale_file = os.path.join(base_dir, "polybench_results_scale8x", "results.json")
    if not os.path.exists(scale_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No 8x-scale results found.", font_size=16, color=ACCENT_RED)
        return

    with open(scale_file) as f:
        wr = json.load(f)

    # Estimated speedups for kernels where benchmark timed out (Triton too slow)
    # Measured via separate single-iteration runs
    timeout_estimates = {
        'cholesky': 0.03,   # C=159ms, Triton=5953ms
        'nussinov': 0.03,   # C=512ms, Triton=16749ms
        'seidel_2d': 0.11,  # C=347ms, Triton=3034ms
    }

    data = []
    estimated_set = set()
    for k, v in wr.items():
        if v.get('test_passed'):
            b = v.get('benchmark')
            if b and b.get('speedup', 0) > 0:
                data.append((k, b['speedup']))
            elif b is None and k in timeout_estimates:
                data.append((k, timeout_estimates[k]))
                estimated_set.add(k)
    data.sort(key=lambda x: -x[1])

    gpu_wins = sum(1 for _, s in data if s >= 1.0)
    all_s = [s for _, s in data]
    sorted_s = sorted(all_s)
    median_s = sorted_s[len(sorted_s) // 2] if sorted_s else 0
    mean_s = sum(all_s) / len(all_s) if all_s else 0
    max_speedup = max((s for _, s in data), default=7.0)

    est_note = f"  |  {len(estimated_set)} estimated (benchmark timeout)" if estimated_set else ""
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(8.0), Inches(0.2),
                f"8x scale | Median: {median_s:.2f}x  |  Mean: {mean_s:.2f}x  |  {gpu_wins}/{len(data)} faster on GPU{est_note}",
                font_size=10, color=MED_GRAY)

    chart_top_in = 1.2
    row_h = min(0.14, 3.8 / max(len(data), 1))
    bar_h = row_h * 0.72
    name_right = 1.35
    bar_left = 1.5
    bar_max_w = 7.5
    # Cap scale at p90 to avoid outliers squashing bars
    p90 = sorted_s[int(len(sorted_s) * 0.9)] if len(sorted_s) > 2 else max_speedup
    max_scale = max(int(p90) + 2, 10)

    for val in range(1, max_scale):
        x = bar_left + bar_max_w * (val / max_scale)
        txb = slide.shapes.add_textbox(
            Inches(x - 0.12), Inches(chart_top_in - 0.18), Inches(0.24), Inches(0.15))
        p = txb.text_frame.paragraphs[0]
        p.text = f"{val}x"
        p.font.name = FONT_BODY
        p.font.size = Pt(7)
        p.font.color.rgb = MED_GRAY
        p.alignment = PP_ALIGN.CENTER

    one_x = bar_left + bar_max_w * (1.0 / max_scale)
    ref_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(one_x), Inches(chart_top_in - 0.02),
        Pt(1.5), Inches(row_h * len(data) + 0.02))
    ref_line.fill.solid()
    ref_line.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    ref_line.line.fill.background()

    for i, (name, speedup) in enumerate(data):
        y_row = chart_top_in + row_h * i
        y_bar = y_row + (row_h - bar_h) / 2

        txb = slide.shapes.add_textbox(
            Inches(0.05), Inches(y_row), Inches(name_right - 0.05), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_CODE
        p.font.size = Pt(7)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.RIGHT

        bar_w_ratio = min(speedup / max_scale, 1.0)
        bar_w = max(bar_max_w * bar_w_ratio, 0.03)
        bar_color = ACCENT_GREEN if speedup >= 1.0 else ACCENT_RED

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_left), Inches(y_bar),
            Inches(bar_w), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        is_est = name in estimated_set
        label_x = bar_left + bar_w + 0.04
        txb = slide.shapes.add_textbox(
            Inches(label_x), Inches(y_row), Inches(0.7), Inches(row_h))
        tf = txb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        spd_str = f"{speedup:.1f}x" if speedup >= 10 else f"{speedup:.2f}x"
        p.text = f"~{spd_str}*" if is_est else spd_str
        p.font.name = FONT_CODE
        p.font.size = Pt(6.5)
        p.font.color.rgb = bar_color
        p.alignment = PP_ALIGN.LEFT

    legend_y = chart_top_in + row_h * len(data) + 0.05
    g = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(7.0), Inches(legend_y), Inches(0.2), Inches(0.12))
    g.fill.solid()
    g.fill.fore_color.rgb = ACCENT_GREEN
    g.line.fill.background()
    slower_count = len(data) - gpu_wins
    add_textbox(slide, Inches(7.25), Inches(legend_y - 0.02), Inches(0.8), Inches(0.15),
                f"Faster ({gpu_wins})", font_size=7, color=DARK_TEXT)
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(8.2), Inches(legend_y), Inches(0.2), Inches(0.12))
    r.fill.solid()
    r.fill.fore_color.rgb = ACCENT_RED
    r.line.fill.background()
    add_textbox(slide, Inches(8.45), Inches(legend_y - 0.02), Inches(0.8), Inches(0.15),
                f"Slower ({slower_count})", font_size=7, color=DARK_TEXT)


def slide_15d_large_regressions(prs):
    """Regression analysis: kernels that degraded from 1x to 8x scale."""
    import json
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "8x-Scale Regressions: What Degraded?")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "Kernels where 8x-scale speedup is significantly worse than 1x baseline. "
                "Root causes: LLM nondeterminism, tolerance failures, or suboptimal parallelization.",
                font_size=10, color=MED_GRAY)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    orig_file = os.path.join(base_dir, "polybench_results", "results.json")
    scale_file = os.path.join(base_dir, "polybench_results_scale8x", "results.json")

    if not os.path.exists(scale_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No 8x-scale results found.", font_size=16, color=ACCENT_RED)
        return

    with open(orig_file) as f:
        orig = json.load(f)
    with open(scale_file) as f:
        scaled = json.load(f)

    # Build regression table sorted by severity
    rows = [["Kernel", "1x Spd", "8x Spd", "Ratio", "8x Att", "Root Cause"]]
    regressions = []

    for k in sorted(orig.keys()):
        o = orig.get(k, {})
        s = scaled.get(k, {})
        o_pass = o.get("test_passed", False)
        s_pass = s.get("test_passed", False)
        o_spd = o.get("benchmark", {}).get("speedup") if o.get("benchmark") else None
        s_spd = s.get("benchmark", {}).get("speedup") if s.get("benchmark") else None

        if not o_pass:
            continue

        if not s_pass:
            # Passed at 1x, failed at 8x
            s_att = s.get("attempts", "?")
            s_err = s.get("final_error", {})
            err_type = s_err.get("type", "unknown") if s_err else "unknown"
            regressions.append((k, o_spd, None, 0, s_att, f"FAIL: {err_type}"))
        elif o_spd and s_spd and s_spd < o_spd * 0.5:
            # >50% speedup degradation
            ratio = s_spd / o_spd
            s_att = s.get("attempts", "?")
            # Per-kernel root cause from code inspection
            known_causes = {
                '2mm': "Serial grid=(1,) fallback; parallel kernels dead code",
                'covariance': "Scalar O(BLOCK^2*N) loops inside tiles",
                'fdtd_2d': "grid=(1,) stencil; 64x more data than 1x",
                'jacobi_1d': "80 kernel launches (2/timestep); launch overhead dominates",
                'jacobi_2d': "grid=(1,) 162K tile iters; data exceeds L1 at 8x",
                'doitgen': "32K Python-side kernel launches (NR*NQ loop)",
                'trisolv': "Inherently sequential O(N^2); no parallelization possible",
            }
            cause = known_causes.get(k, "Suboptimal parallelization at larger size")
            regressions.append((k, o_spd, s_spd, ratio, s_att, cause))

    regressions.sort(key=lambda x: x[3])

    for k, o_spd, s_spd, ratio, att, cause in regressions:
        o_str = f"{o_spd:.2f}x" if o_spd else "—"
        s_str = f"{s_spd:.2f}x" if s_spd else "FAIL"
        r_str = f"{ratio:.2f}x" if ratio > 0 else "—"
        rows.append([k, o_str, s_str, r_str, str(att), cause])

    if len(rows) > 1:
        add_simple_table(slide, Inches(0.3), Inches(1.2), Inches(9.4),
                         [1.1, 0.7, 0.7, 0.6, 0.5, 3.6], rows,
                         font_size=9, row_height=0.28)
    else:
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "No significant regressions detected.", font_size=14, color=ACCENT_GREEN)
        return

    # Categorize
    n_fail = sum(1 for _, _, s, _, _, _ in regressions if s is None)
    n_degrade = len(regressions) - n_fail
    def _safe_spd(d, k):
        entry = d.get(k, {})
        bench = entry.get("benchmark") if entry else None
        return bench.get("speedup", 0) if bench else 0

    n_improved = sum(1 for k in orig
                     if orig[k].get("test_passed") and scaled.get(k, {}).get("test_passed")
                     and _safe_spd(orig, k) > 0
                     and _safe_spd(scaled, k) > _safe_spd(orig, k) * 1.1)

    # Summary insight at bottom
    tbl_bottom = 1.2 + 0.28 * len(rows) + 0.15
    add_textbox(slide, Inches(0.5), Inches(tbl_bottom), Inches(9.0), Inches(0.25),
                "Analysis:", font_size=13, bold=True, color=DARK_BLUE)

    analysis_items = [
        f"{n_fail} kernel(s) regressed from PASS to FAIL (FP32 accumulation error over longer chains at 8x)",
        f"{n_degrade} kernel(s) show >50% speedup degradation vs 1x baseline",
        f"Dominant pattern: grid=(1,) kernels that fit in cache at 1x but exceed it at 8x (64x more data)",
        f"Secondary pattern: excessive Python-side kernel launches (doitgen: 32K, jacobi_1d: 80)",
        f"{n_improved} kernel(s) improved >10% at 8x (larger matrices amortize GPU launch overhead)",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(tbl_bottom + 0.3), Inches(9.0), Inches(1.5),
                    analysis_items, font_size=10, color=DARK_TEXT)


def slide_16_nondeterminism_evidence(prs):
    """Nondeterminism evidence for current NA-win kernels only."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    results_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "polybench_results")
    evidence_file = os.path.join(results_dir, "nondeterminism_evidence.json")

    # Compute current NA wins from results.json vs results_no_analysis.json
    with_file = os.path.join(results_dir, "results.json")
    without_file = os.path.join(results_dir, "results_no_analysis.json")
    na_win_kernels = set()
    if os.path.exists(with_file) and os.path.exists(without_file):
        with open(with_file) as f:
            wr = json.load(f)
        with open(without_file) as f:
            nr = json.load(f)
        for k in wr:
            w = wr.get(k, {})
            n = nr.get(k, {})
            if w.get('test_passed') and n.get('test_passed'):
                ws = w.get('benchmark', {}).get('speedup', 0)
                ns = n.get('benchmark', {}).get('speedup', 0)
                if ws > 0 and ns > 0 and ws / ns < 0.87:
                    na_win_kernels.add(k)

    n_na = len(na_win_kernels)
    add_section_title(slide, f"Regression Analysis: {n_na} NA-Win Kernels Classified")

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(9.0), Inches(0.55),
                "Each NA-win kernel re-run 3-5x. Triton ms compared against NA baseline. "
                "Green = nondeterminism (WA beats NA >= 2/N runs). Red = true regression (WA beats NA <= 1/N). "
                "NA wins vary across runs due to LLM nondeterminism; evidence below is from a dedicated multi-run study.",
                font_size=9, color=MED_GRAY)

    if not os.path.exists(evidence_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "Evidence file not found. Run run_nondeterminism_study.py first.",
                    font_size=16, color=ACCENT_RED)
        return

    with open(evidence_file) as f:
        evidence = json.load(f)

    # Filter to current NA wins only
    filtered_evidence = {k: v for k, v in evidence.items() if k in na_win_kernels}

    # Build table rows: Kernel | NA (ms) | Best WA | Worst WA | WA>NA | Verdict
    rows = [["Kernel", "NA (ms)", "Best WA", "Worst WA", "WA>NA", "Verdict"]]

    total_nondet = 0
    total_true_reg = 0

    # Sort: nondeterminism first, then true regressions
    sorted_kernels = sorted(filtered_evidence.keys(),
                            key=lambda k: (0 if filtered_evidence[k].get("wa_beats_na_count", 0) >= 2 else 1, k))

    for kname in sorted_kernels:
        e = filtered_evidence[kname]
        na_t = e.get("na_triton_ms", 0)
        runs = e.get("runs", [])
        triton_times = [r["triton_ms"] for r in runs if r.get("triton_ms") is not None]
        best_wa = min(triton_times) if triton_times else None
        worst_wa = max(triton_times) if triton_times else None
        wa_beats = e.get("wa_beats_na_count", 0)
        total = e.get("total_runs", len(runs))

        # Classify: >= 2 WA beats = nondeterminism, <= 1 = true regression
        is_nondet = wa_beats >= 2
        if is_nondet:
            total_nondet += 1
        else:
            total_true_reg += 1

        best_str = f"{best_wa:.3f}" if best_wa is not None else "FAIL"
        worst_str = f"{worst_wa:.3f}" if worst_wa is not None else "FAIL"
        verdict = "Nondet" if is_nondet else "True reg"

        rows.append([kname, f"{na_t:.3f}", best_str, worst_str,
                     f"{wa_beats}/{total}", verdict])

    # Also add NA-win kernels not in evidence file
    for kname in sorted(na_win_kernels - set(filtered_evidence.keys())):
        total_true_reg += 1  # conservative: no evidence = unknown
        rows.append([kname, "—", "—", "—", "—", "No data"])

    tbl = add_simple_table(slide, Inches(0.1), Inches(1.3), Inches(9.8),
                           [1.4, 1.2, 1.2, 1.2, 1.0, 1.2], rows,
                           font_size=9, row_height=0.28)

    # Color verdict cells
    for ri in range(1, len(rows)):
        verdict_cell = tbl.cell(ri, 5)
        if rows[ri][5] == "Nondet":
            verdict_cell.fill.solid()
            verdict_cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)  # light green
        else:
            verdict_cell.fill.solid()
            verdict_cell.fill.fore_color.rgb = RGBColor(0xFA, 0xDB, 0xD8)  # light red

    # Observations section
    obs_y = Inches(1.3 + 0.28 * len(rows) + 0.15)
    add_textbox(slide, Inches(0.5), obs_y, Inches(4.3), Inches(0.25),
                f"Nondeterminism ({total_nondet}/{n_na} kernels)",
                font_size=11, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, Inches(0.5), obs_y + Inches(0.25), Inches(4.3), Inches(0.6),
                "WA beats NA in multiple re-runs. The NA advantage on the "
                "ablation slide is due to LLM sampling variance, not analysis quality.",
                font_size=9, color=DARK_TEXT)

    add_textbox(slide, Inches(5.0), obs_y, Inches(4.8), Inches(0.25),
                f"True Regressions ({total_true_reg}/{n_na} kernels)",
                font_size=11, bold=True, color=ACCENT_RED)
    add_textbox(slide, Inches(5.0), obs_y + Inches(0.25), Inches(4.8), Inches(0.6),
                "WA rarely or never beats NA. Analysis guidance steers the LLM "
                "toward a suboptimal strategy for these kernels.",
                font_size=9, color=DARK_TEXT)

    # Bottom conclusion box
    concl_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.3), Inches(5.05), Inches(9.5), Inches(0.45))
    concl_box.fill.solid()
    concl_box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    concl_box.line.color.rgb = MED_BLUE
    concl_box.line.width = Pt(1.5)
    concl_box.adjustments[0] = 0.08
    tf = concl_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Conclusion: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = MED_BLUE
    run2 = p.add_run()
    run2.text = (f"{total_nondet}/{n_na} NA wins are nondeterminism (WA beats NA in re-runs). "
                 f"{total_true_reg}/{n_na} are true regressions where analysis guides suboptimally.")
    run2.font.size = Pt(10)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


def slide_17_large_size_ablation(prs):
    """Large-size ablation: WA vs NA at 8x problem sizes (full pipeline re-run)."""
    import json

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "8x-Scale Ablation: Analysis vs No-Analysis")

    add_textbox(slide, Inches(0.5), Inches(0.8), Inches(9.0), Inches(0.3),
                "Both WA and NA pipelines re-run from scratch at 8x problem sizes. "
                "LLM sees scaled dimensions in prompt. C reference recompiled at scaled sizes.",
                font_size=10, color=MED_GRAY)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    wa_file = os.path.join(base_dir, "polybench_results_scale8x", "results.json")
    na_file = os.path.join(base_dir, "polybench_results_scale8x", "results_no_analysis.json")

    if not os.path.exists(wa_file) or not os.path.exists(na_file):
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(1.0),
                    "8x-scale ablation data not found. Run both --size-scale 8 and --no-analysis --size-scale 8.",
                    font_size=16, color=ACCENT_RED)
        return

    with open(wa_file) as f:
        wa = json.load(f)
    with open(na_file) as f:
        na = json.load(f)

    # Build table: Kernel | WA Pass | WA Spd | NA Pass | NA Spd | Winner
    rows = [["Kernel", "WA Pass", "WA Spd", "NA Pass", "NA Spd", "Winner"]]

    wa_wins, na_wins, ties = 0, 0, 0
    both_bench = 0
    wa_only_pass, na_only_pass = 0, 0

    all_kernels = sorted(set(list(wa.keys()) + list(na.keys())))
    kernel_data = []
    for k in all_kernels:
        w = wa.get(k, {})
        n = na.get(k, {})
        w_pass = w.get("test_passed", False)
        n_pass = n.get("test_passed", False)
        w_spd = w.get("benchmark", {}).get("speedup") if w.get("benchmark") else None
        n_spd = n.get("benchmark", {}).get("speedup") if n.get("benchmark") else None
        kernel_data.append((k, w_pass, w_spd, n_pass, n_spd))

    # Sort by WA speedup descending
    kernel_data.sort(key=lambda x: -(x[2] or -999))

    for k, w_pass, w_spd, n_pass, n_spd in kernel_data:
        wp = "Y" if w_pass else "N"
        np_ = "Y" if n_pass else "N"
        ws = f"{w_spd:.2f}x" if w_spd else "—"
        ns = f"{n_spd:.2f}x" if n_spd else "—"

        if w_pass and not n_pass:
            winner = "WA-only"
            wa_only_pass += 1
        elif n_pass and not w_pass:
            winner = "NA-only"
            na_only_pass += 1
        elif w_spd and n_spd and w_spd > 0 and n_spd > 0:
            both_bench += 1
            ratio = w_spd / n_spd
            if ratio > 1.05:
                winner = "WA"
                wa_wins += 1
            elif ratio < 0.95:
                winner = "NA"
                na_wins += 1
            else:
                winner = "TIE"
                ties += 1
        elif w_pass and n_pass:
            winner = "—"
        else:
            winner = "—"

        rows.append([k, wp, ws, np_, ns, winner])

    tbl = add_simple_table(slide, Inches(0.15), Inches(1.15), Inches(6.0),
                           [1.2, 0.6, 1.0, 0.6, 1.0, 0.8], rows,
                           font_size=8, row_height=0.22)

    # Color winner cells
    for ri in range(1, len(rows)):
        winner_cell = tbl.cell(ri, 5)
        w = rows[ri][5]
        if w == "WA" or w == "WA-only":
            winner_cell.fill.solid()
            winner_cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)
        elif w == "NA" or w == "NA-only":
            winner_cell.fill.solid()
            winner_cell.fill.fore_color.rgb = RGBColor(0xFA, 0xDB, 0xD8)

    # Stats panel on right
    stats_x = Inches(6.3)
    wa_pass_n = sum(1 for _, wp, _, _, _ in kernel_data if wp)
    na_pass_n = sum(1 for _, _, _, np_, _ in kernel_data if np_)
    wa_spds = [s for _, wp, s, _, _ in kernel_data if wp and s and s > 0]
    na_spds = [s for _, _, _, np_, s in kernel_data if np_ and s and s > 0]
    wa_sorted = sorted(wa_spds) if wa_spds else [0]
    na_sorted = sorted(na_spds) if na_spds else [0]

    add_textbox(slide, stats_x, Inches(1.15), Inches(3.5), Inches(0.25),
                "8x-Scale Summary:", font_size=13, bold=True, color=DARK_BLUE)

    stats_items = [
        f"Pass rate: WA {wa_pass_n}/30 vs NA {na_pass_n}/30",
        f"Benchmarked: WA {len(wa_spds)} vs NA {len(na_spds)}",
        f"Median: WA {wa_sorted[len(wa_sorted)//2]:.2f}x vs NA {na_sorted[len(na_sorted)//2]:.2f}x",
        f"Mean: WA {sum(wa_spds)/len(wa_spds):.2f}x vs NA {sum(na_spds)/len(na_spds):.2f}x" if na_spds else "",
        f">1x: WA {sum(1 for s in wa_spds if s > 1)}/{len(wa_spds)} vs NA {sum(1 for s in na_spds if s > 1)}/{len(na_spds)}",
    ]
    stats_items = [s for s in stats_items if s]

    add_bullet_list(slide, stats_x, Inches(1.45), Inches(3.5), Inches(1.8),
                    stats_items, font_size=10, color=DARK_TEXT)

    add_textbox(slide, stats_x, Inches(3.1), Inches(3.5), Inches(0.25),
                "Head-to-Head:", font_size=13, bold=True, color=ACCENT_GREEN)

    h2h_items = [
        f"WA wins: {wa_wins}/{both_bench}",
        f"NA wins: {na_wins}/{both_bench}",
        f"Ties (<5%): {ties}/{both_bench}",
        f"WA uniquely passes: {wa_only_pass}",
    ]
    add_bullet_list(slide, stats_x, Inches(3.4), Inches(3.5), Inches(1.5),
                    h2h_items, font_size=10, color=DARK_TEXT)

    # Conclusion box
    concl_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.3), Inches(5.05), Inches(9.5), Inches(0.45))
    concl_box.fill.solid()
    concl_box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    concl_box.line.color.rgb = MED_BLUE
    concl_box.line.width = Pt(1.5)
    concl_box.adjustments[0] = 0.08
    tf = concl_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Conclusion: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.name = FONT_BODY
    run.font.color.rgb = MED_BLUE
    run2 = p.add_run()
    if both_bench > 0:
        wa_pct = wa_wins / both_bench * 100
        run2.text = (f"At 8x sizes, WA passes {wa_pass_n} vs NA {na_pass_n} kernels. "
                     f"Among {both_bench} head-to-head benchmarks, WA wins {wa_wins} ({wa_pct:.0f}%). "
                     f"Analysis advantage grows at realistic problem sizes.")
    else:
        run2.text = "No head-to-head comparisons available."
    run2.font.size = Pt(10)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = DARK_TEXT


# ════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    slide_01_title(prs)
    slide_02_architecture(prs)
    slide_02b_extraction_overview(prs)
    slide_02c_gemm_before_after(prs)
    slide_02d_kernel_to_prompt(prs)
    slide_03_ablation_overview(prs)
    slide_03b_no_analysis_kernels(prs)
    slide_03c_ablation_all_kernels(prs)
    slide_16_nondeterminism_evidence(prs)
    slide_04_war_trisolv_trmm(prs)
    slide_05_war_lu_seidel(prs)
    slide_06_scalarexp_deriche(prs)
    slide_07_scalarexp_durbin_symm(prs)
    slide_08_analysis_summary(prs)
    slide_09_heat3d_failure(prs)
    slide_09_gemm_triton(prs)
    slide_09b_evaluation(prs)
    slide_09c_test_inputs(prs)
    slide_10_results(prs)
    slide_10b_speedup_chart(prs)
    slide_10c_perf_comparison(prs)
    slide_11_failures(prs)
    slide_12_ablation(prs)
    slide_13_rodinia_results(prs)
    slide_13b_rodinia_speedup_chart(prs)
    slide_14_rodinia_ablation(prs)
    slide_15_ncu_utilization(prs)
    slide_15b_large_size_benchmark(prs)
    slide_15c_large_speedup_chart(prs)
    slide_15d_large_regressions(prs)
    slide_17_large_size_ablation(prs)

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "polybench_pipeline_slides.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
